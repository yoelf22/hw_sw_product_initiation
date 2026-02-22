#!/usr/bin/env python3
"""Build a LinkedIn carousel PDF and PPTX for Haptic Metronome Bracelet.

Format: 1080x1350 px (4:5 portrait) — optimized for mobile feed.
Uses reportlab for PDF and python-pptx for PPTX generation.
"""

import os
from reportlab.lib.units import mm
from reportlab.lib.colors import HexColor, white
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

_DIR = os.path.dirname(os.path.abspath(__file__))

# -- Page size: 4:5 ratio --
PW = 190 * mm
PH = 237.5 * mm

# -- Colors --
DARK_BG = "#1A1A2E"
CARD_BG = "#22223A"
CARD_BG_ALT = "#1E1E34"
ACCENT_ORANGE = "#FF8C00"
ACCENT_GREEN = "#4EC978"
ACCENT_RED = "#FF4545"
ACCENT_BLUE = "#009BF5"
ACCENT_PURPLE = "#A855F7"
WHITE_HEX = "#FFFFFF"
LIGHT_GRAY = "#BBBBCC"
SOFT_WHITE = "#F0F0F5"

M = 10 * mm  # standard margin


def bg(c, color=DARK_BG):
    c.setFillColor(HexColor(color))
    c.rect(0, 0, PW, PH, fill=1, stroke=0)


def accent_strip(c, color=ACCENT_ORANGE):
    c.setFillColor(HexColor(color))
    c.rect(0, PH - 2 * mm, PW, 2 * mm, fill=1, stroke=0)


def card(c, x, y, w, h, color=CARD_BG):
    """Draw a card. y is from TOP of page."""
    c.setFillColor(HexColor(color))
    c.roundRect(x, PH - y - h, w, h, 3 * mm, fill=1, stroke=0)


def card_flat(c, x, y, w, h, color=CARD_BG):
    c.setFillColor(HexColor(color))
    c.rect(x, PH - y - h, w, h, fill=1, stroke=0)


def txt(c, x, y, text, size=14, color=WHITE_HEX, bold=False, align="left", max_w=None):
    """Draw text. y is from TOP of page."""
    c.setFillColor(HexColor(color))
    font = "Helvetica-Bold" if bold else "Helvetica"
    c.setFont(font, size)
    ry = PH - y - size * 0.35
    if align == "center" and max_w:
        tw = c.stringWidth(text, font, size)
        c.drawString(x + (max_w - tw) / 2, ry, text)
    elif align == "right" and max_w:
        tw = c.stringWidth(text, font, size)
        c.drawString(x + max_w - tw, ry, text)
    else:
        c.drawString(x, ry, text)


def txt_wrap(c, x, y, text, size=11, color=WHITE_HEX, bold=False, line_h=None, max_w=None):
    """Draw wrapped text. Returns y after last line."""
    if line_h is None:
        line_h = size * 1.4
    font = "Helvetica-Bold" if bold else "Helvetica"
    c.setFillColor(HexColor(color))
    c.setFont(font, size)
    if max_w is None:
        max_w = PW - 2 * M
    words = text.split()
    lines = []
    current = ""
    for w in words:
        test = current + (" " if current else "") + w
        if c.stringWidth(test, font, size) > max_w:
            if current:
                lines.append(current)
            current = w
        else:
            current = test
    if current:
        lines.append(current)
    for line in lines:
        c.drawString(x, PH - y - size * 0.35, line)
        y += line_h
    return y


def bar(c, x, y, w, h, color):
    c.setFillColor(HexColor(color))
    c.rect(x, PH - y - h, w, h, fill=1, stroke=0)


def circle_num(c, x, y, num, color):
    r = 4 * mm
    c.setFillColor(HexColor(color))
    c.circle(x + r, PH - y - r, r, fill=1, stroke=0)
    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 12)
    tw = c.stringWidth(str(num), "Helvetica-Bold", 12)
    c.drawString(x + r - tw / 2, PH - y - r - 4, str(num))


def footer(c, page, total):
    c.setFillColor(HexColor(LIGHT_GRAY))
    c.setFont("Helvetica", 8)
    label = f"{page}/{total}"
    tw = c.stringWidth(label, "Helvetica", 8)
    c.drawString(PW - M - tw, 4 * mm, label)


TOTAL_PAGES = 8
output = os.path.join(_DIR, "Haptic_Metronome_Bracelet_Carousel.pdf")
c = canvas.Canvas(output, pagesize=(PW, PH))

# ================================================================
# PAGE 1: Title
# ================================================================
bg(c)
accent_strip(c, ACCENT_PURPLE)

txt(c, M, 14 * mm, "Haptic Metronome", size=38, color=WHITE_HEX, bold=True)
txt(c, M, 28 * mm, "Bracelet", size=38, color=WHITE_HEX, bold=True)

bar(c, M, 37 * mm, 35 * mm, 1 * mm, ACCENT_PURPLE)

txt_wrap(c, M, 43 * mm,
         "Feel the beat. Hear nothing.",
         size=18, color=ACCENT_GREEN, max_w=PW - 2 * M)

txt_wrap(c, M, 56 * mm,
         "A wrist-worn vibrotactile metronome for musicians who need "
         "silent, precise time -- practice, rehearsal, and stage.",
         size=11, color=LIGHT_GRAY, max_w=PW - 2 * M)

# Image
img_path = os.path.join(_DIR, "cross_section_illustration_haptic_metronome.png")
if os.path.exists(img_path):
    img_top_y = 75 * mm
    img_max_w = PW - 2 * M
    img_max_h = 135 * mm
    c.drawImage(ImageReader(img_path),
                M, PH - img_top_y - img_max_h,
                width=img_max_w, height=img_max_h,
                preserveAspectRatio=True, anchor='sw', mask='auto')

# Bottom bar
card_flat(c, 0, PH - 8 * mm, PW, 8 * mm, CARD_BG)
txt(c, M, PH - 5 * mm, "Product Overview  |  Concept Stage  |  2026",
    size=9, color=LIGHT_GRAY)

footer(c, 1, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 2: The Problem
# ================================================================
bg(c)
accent_strip(c, ACCENT_RED)

txt(c, M, 12 * mm, "The Problem", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_RED)

problems = [
    ("Audible clicks are unusable on stage",
     "Click tracks bleed into microphones, in-ear monitors isolate "
     "from the room, and audible metronomes distract the audience.",
     ACCENT_RED),
    ("Phone apps are imprecise",
     "Android audio latency: 5-40 ms. iOS: 5-15 ms. OS interrupts "
     "and background tasks cause timing jitter. Phone vibration motors "
     "are too slow and vague for musical time.",
     ACCENT_ORANGE),
    ("Existing solutions compromise",
     "Earpiece click tracks occupy an ear. Visual metronomes require "
     "looking away from the music. There is no silent, precise, "
     "hands-free time reference that leaves both ears and eyes open.",
     ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(problems):
    y_top = 30 * mm + i * 55 * mm
    card(c, M, y_top, PW - 2 * M, 49 * mm, CARD_BG)
    bar(c, M, y_top, PW - 2 * M, 1.5 * mm, color)
    txt(c, M + 5 * mm, y_top + 8 * mm, title, size=16, color=color, bold=True)
    txt_wrap(c, M + 5 * mm, y_top + 20 * mm, desc,
             size=11, color=LIGHT_GRAY, max_w=PW - 2 * M - 10 * mm, line_h=14)

# Target users
card_flat(c, M, 200 * mm, PW - 2 * M, 30 * mm, CARD_BG)
txt(c, M + 4 * mm, 204 * mm, "TARGET USERS", size=10, color=ACCENT_PURPLE, bold=True)
txt(c, M + 4 * mm, 213 * mm, "Guitarists  |  Pianists  |  String players",
    size=12, color=WHITE_HEX, bold=True)
txt(c, M + 4 * mm, 222 * mm, "Wind & brass  |  Vocalists  |  Conductors",
    size=12, color=WHITE_HEX, bold=True)

footer(c, 2, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 3: How It Works
# ================================================================
bg(c)
accent_strip(c, ACCENT_GREEN)

txt(c, M, 12 * mm, "How It Works", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_GREEN)

steps = [
    ("Set", "Open app, set BPM and time signature",
     "App pushes config over BLE to the bracelet. "
     "Save presets and setlists for quick recall.",
     ACCENT_BLUE),
    ("Tap", "Tap the bracelet to start",
     "Accelerometer detects a fingertip tap on the pod. "
     "Haptic pulses begin immediately. LED flashes on beat "
     "for the first 4 bars as a sanity check, then goes dark.",
     ACCENT_GREEN),
    ("Feel", "Feel the beat on your wrist",
     "Single pulse = normal beat (~15 ms crisp tap). "
     "Dual-pulse = downbeat (\"da-dum\", two taps 40 ms apart). "
     "You always know where beat 1 is.",
     ACCENT_PURPLE),
    ("Stop", "Double-tap to stop",
     "Accelerometer detects the double-tap gesture. "
     "Pulses stop. Device goes to low-power idle, "
     "ready for the next tap.",
     ACCENT_ORANGE),
]

for i, (title, subtitle, desc, color) in enumerate(steps):
    y_top = 28 * mm + i * 50 * mm
    card(c, M, y_top, PW - 2 * M, 44 * mm, CARD_BG)
    circle_num(c, M + 4 * mm, y_top + 4 * mm, i + 1, color)
    txt(c, M + 16 * mm, y_top + 6 * mm, title, size=18, color=color, bold=True)
    txt(c, M + 16 * mm, y_top + 16 * mm, subtitle, size=11, color=WHITE_HEX, bold=True)
    txt_wrap(c, M + 6 * mm, y_top + 26 * mm, desc,
             size=10, color=LIGHT_GRAY, max_w=PW - 2 * M - 12 * mm, line_h=13)

# Bottom note
card_flat(c, M, 232 * mm, PW - 2 * M, 4 * mm, CARD_BG)

footer(c, 3, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 4: Architecture
# ================================================================
bg(c)
accent_strip(c, ACCENT_BLUE)

txt(c, M, 12 * mm, "Architecture", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_BLUE)

# Signal chain
card(c, M, 28 * mm, PW - 2 * M, 28 * mm, CARD_BG)
txt(c, M + 4 * mm, 32 * mm, "Beat Signal Chain", size=11, color=ACCENT_BLUE, bold=True)
txt(c, M + 4 * mm, 41 * mm, "Timer ISR  -->  I2C cmd  -->  DRV2605L  -->  LRA motor",
    size=12, color=WHITE_HEX, bold=True)
txt(c, M + 4 * mm, 49 * mm, "< 100 us jitter      waveform select     shaped drive     crisp tap on wrist",
    size=8, color=LIGHT_GRAY)

# Gesture chain
card(c, M, 60 * mm, PW - 2 * M, 22 * mm, CARD_BG)
txt(c, M + 4 * mm, 64 * mm, "Gesture Chain", size=11, color=ACCENT_GREEN, bold=True)
txt(c, M + 4 * mm, 72 * mm, "Tap on pod  -->  LIS2DH12 interrupt  -->  FW validation  -->  Start/Stop",
    size=10, color=WHITE_HEX, bold=True)

# Subsystems
bar(c, M, 88 * mm, PW - 2 * M, 1 * mm, ACCENT_BLUE)
txt(c, M, 92 * mm, "SUBSYSTEMS", size=11, color=ACCENT_BLUE, bold=True)

subsystems = [
    ("MCU (nRF52832)", "Timing engine, BLE 5.0, gesture detection, haptic control", ACCENT_GREEN),
    ("Haptic (DRV2605L + LRA)", "Crisp vibrotactile pulses, auto-resonance tracking", ACCENT_PURPLE),
    ("Accelerometer (LIS2DH12)", "Hardware tap/double-tap detection, Z-axis validation", ACCENT_GREEN),
    ("BLE 5.0 (GATT)", "App config, presets, battery status, OTA updates", ACCENT_BLUE),
    ("Power (150mAh LiPo)", "USB-C charging, 30+ hrs continuous play", ACCENT_ORANGE),
    ("LED (1x green)", "Sanity check: first 4 bars only, then dark", ACCENT_ORANGE),
]

for i, (name, desc, color) in enumerate(subsystems):
    y_top = 98 * mm + i * 17 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 14 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 4 * mm, name, size=11, color=WHITE_HEX, bold=True)
    txt(c, M + 4 * mm, y_top + 11 * mm, desc, size=8, color=LIGHT_GRAY)
    bar(c, M, y_top, 1.5 * mm, 14 * mm, color)

# App as primary UI
card(c, M, 202 * mm, PW - 2 * M, 28 * mm, CARD_BG_ALT)
txt(c, M + 4 * mm, 206 * mm, "Thin actuator, thick app", size=13, color=ACCENT_PURPLE, bold=True)
txt_wrap(c, M + 4 * mm, 216 * mm,
         "The bracelet handles timing and haptics. The app handles everything "
         "else: BPM, meter, accent patterns, presets, setlists, practice logs. "
         "No physical controls for config -- tap/double-tap for start/stop only.",
         size=9, color=LIGHT_GRAY, max_w=PW - 2 * M - 8 * mm, line_h=12)

footer(c, 4, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 5: Downbeat Differentiation
# ================================================================
bg(c)
accent_strip(c, ACCENT_PURPLE)

txt(c, M, 12 * mm, "Feeling the Downbeat", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_PURPLE)

txt_wrap(c, M, 28 * mm,
         "How do you know where beat 1 is when you can only feel vibration?",
         size=13, color=LIGHT_GRAY, max_w=PW - 2 * M)

# Normal beat
card(c, M, 42 * mm, PW - 2 * M, 45 * mm, CARD_BG)
bar(c, M, 42 * mm, PW - 2 * M, 1.5 * mm, ACCENT_GREEN)
txt(c, M + 5 * mm, 48 * mm, "Normal Beat (2, 3, 4...)", size=16, color=ACCENT_GREEN, bold=True)
txt(c, M + 5 * mm, 60 * mm, "Single pulse: one crisp tap", size=13, color=WHITE_HEX)
txt(c, M + 5 * mm, 71 * mm, "|--15ms--|", size=11, color=LIGHT_GRAY)
# Visual pulse representation
bar(c, M + 5 * mm, 78 * mm, 15 * mm, 3 * mm, ACCENT_GREEN)

# Downbeat
card(c, M, 94 * mm, PW - 2 * M, 55 * mm, CARD_BG)
bar(c, M, 94 * mm, PW - 2 * M, 1.5 * mm, ACCENT_PURPLE)
txt(c, M + 5 * mm, 100 * mm, "Downbeat (Beat 1)", size=16, color=ACCENT_PURPLE, bold=True)
txt(c, M + 5 * mm, 112 * mm, "Dual-pulse: \"da-dum\"", size=13, color=WHITE_HEX)
txt(c, M + 5 * mm, 122 * mm, "Two taps, 40 ms apart", size=13, color=WHITE_HEX)
txt(c, M + 5 * mm, 135 * mm, "|--15ms--|----40ms----|--15ms--|", size=11, color=LIGHT_GRAY)
# Visual dual pulse
bar(c, M + 5 * mm, 142 * mm, 15 * mm, 3 * mm, ACCENT_PURPLE)
bar(c, M + 60 * mm, 142 * mm, 15 * mm, 3 * mm, ACCENT_PURPLE)

# Rationale
card(c, M, 158 * mm, PW - 2 * M, 70 * mm, CARD_BG_ALT)
txt(c, M + 5 * mm, 163 * mm, "Why pattern, not intensity?", size=13, color=ACCENT_ORANGE, bold=True)

rationale = [
    ("Amplitude is ambiguous",
     "\"Slightly stronger\" vs. \"normal\" is hard to "
     "distinguish when your arm is moving."),
    ("Pattern is distinct",
     "The brain detects a double-tap as a different event, "
     "not just a louder version of the same event."),
    ("Works at tempo",
     "At 200 BPM (300 ms between beats), the 70 ms dual-pulse "
     "uses 23% of the interval -- tight but clear."),
]

ry = 175 * mm
for title, desc in rationale:
    txt(c, M + 5 * mm, ry, f"- {title}:", size=10, color=WHITE_HEX, bold=True)
    ry = txt_wrap(c, M + 8 * mm, ry + 12, desc,
                  size=9, color=LIGHT_GRAY, max_w=PW - 2 * M - 16 * mm, line_h=12)
    ry += 3 * mm

footer(c, 5, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 6: Constraints & BOM
# ================================================================
bg(c)
accent_strip(c, ACCENT_ORANGE)

txt(c, M, 12 * mm, "Constraints & BOM", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_ORANGE)

constraints = [
    ("Timing jitter", "< 100 us", "Hardware timer ISR, not software loop"),
    ("Haptic rise time", "< 10 ms", "LRA mandatory; ERM too slow above 160 BPM"),
    ("Battery life", "> 8 hrs continuous", "150 mAh LiPo, ~4.9 mA avg during play"),
    ("Pod size", "35 x 25 x 10 mm", "Fits comfortably on wrist during playing"),
    ("Pod weight", "< 20g", "Lighter than most watches"),
    ("Sweat resistance", "IPX4", "Silicone overmold, sealed USB-C flap"),
]

for i, (name, value, note) in enumerate(constraints):
    y_top = 28 * mm + i * 17 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 14 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 3.5 * mm, name, size=11, color=ACCENT_ORANGE, bold=True)
    txt(c, M + 4 * mm, y_top + 10 * mm, note, size=8, color=LIGHT_GRAY)
    txt(c, M + 4 * mm, y_top + 3.5 * mm, value, size=11, color=WHITE_HEX, bold=True,
        align="right", max_w=PW - 2 * M - 8 * mm)

# BOM
bar(c, M, 133 * mm, PW - 2 * M, 1 * mm, ACCENT_GREEN)
txt(c, M, 137 * mm, "BOM ESTIMATE (1k units)", size=11, color=ACCENT_GREEN, bold=True)

bom = [
    ("nRF52832 MCU", "$2.80"),
    ("DRV2605L haptic driver", "$1.50"),
    ("LRA motor (8mm)", "$1.50"),
    ("LIS2DH12 accelerometer", "$0.80"),
    ("LiPo 150mAh + protection", "$1.60"),
    ("Crystals + LED + passives", "$0.70"),
    ("USB-C + charge IC + LDO", "$0.95"),
    ("PCB (28x18mm, 4-layer)", "$0.80"),
    ("Pod enclosure (silicone overmold)", "$2.50"),
    ("Silicone wristband + lugs", "$1.20"),
    ("Packaging + USB-C cable", "$1.20"),
    ("Assembly + test", "$2.50"),
]

for i, (item, cost) in enumerate(bom):
    y_top = 144 * mm + i * 7 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card_flat(c, M + 2 * mm, y_top, PW - 2 * M - 4 * mm, 6 * mm, bg_c)
    txt(c, M + 6 * mm, y_top + 1.5 * mm, item, size=8, color=SOFT_WHITE)
    txt(c, M + 6 * mm, y_top + 1.5 * mm, cost, size=8, color=WHITE_HEX, bold=True,
        align="right", max_w=PW - 2 * M - 16 * mm)

# Total
card_flat(c, M + 2 * mm, 228 * mm, PW - 2 * M - 4 * mm, 7 * mm, ACCENT_ORANGE)
txt(c, M + 6 * mm, 230 * mm, "Total COGS", size=10, color=WHITE_HEX, bold=True)
txt(c, M + 6 * mm, 230 * mm, "~$18.05  (5k: ~$14.50)", size=10, color=WHITE_HEX, bold=True,
    align="right", max_w=PW - 2 * M - 16 * mm)

footer(c, 6, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 7: Hardest Problems
# ================================================================
bg(c)
accent_strip(c, ACCENT_RED)

txt(c, M, 12 * mm, "Hardest Problems", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_RED)

problems = [
    ("Haptic perceptibility during playing",
     "Can the LRA pulse (~1.5g) be felt on the wrist while "
     "strumming guitar, bowing violin, or playing piano expressively? "
     "Arm motion creates competing vibrations. Tight skin contact "
     "and sharp LRA rise time are the primary mitigations. "
     "Must prototype and test with real musicians."),
    ("Tap gesture vs. playing motion",
     "The accelerometer must reliably detect a deliberate fingertip "
     "tap on the pod while rejecting strumming, bowing, and arm "
     "movement. Two-stage detection: hardware threshold trigger, then "
     "firmware Z-axis dominance validation. Target: < 1% false "
     "trigger, < 5% missed tap."),
    ("Downbeat perception at tempo",
     "Can musicians distinguish the dual-pulse downbeat from a "
     "single pulse at 120, 160, and 200 BPM? At 200 BPM the "
     "70 ms dual-pulse occupies 23% of the beat interval. The 40 ms "
     "gap must be above tactile temporal resolution (~5-10 ms). "
     "Needs testing with 5+ musicians across instruments."),
]

for i, (title, desc) in enumerate(problems):
    y_top = 30 * mm + i * 62 * mm
    card(c, M, y_top, PW - 2 * M, 56 * mm, CARD_BG)
    circle_num(c, M + 4 * mm, y_top + 4 * mm, i + 1, ACCENT_RED)
    txt(c, M + 16 * mm, y_top + 6 * mm, title, size=13, color=WHITE_HEX, bold=True)
    txt_wrap(c, M + 6 * mm, y_top + 18 * mm, desc,
             size=10, color=LIGHT_GRAY, max_w=PW - 2 * M - 12 * mm, line_h=13)

# Bottom
card_flat(c, M, 218 * mm, PW - 2 * M, 16 * mm, CARD_BG)
txt_wrap(c, M + 4 * mm, 222 * mm,
         "All three require prototype validation with real musicians. "
         "No amount of simulation replaces putting this on a wrist.",
         size=11, color=ACCENT_ORANGE, max_w=PW - 2 * M - 8 * mm, line_h=14)

footer(c, 7, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 8: Gate Result & Open Items
# ================================================================
bg(c)
accent_strip(c, ACCENT_GREEN)

txt(c, M, 12 * mm, "Gate Result & Next", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_GREEN)

# Gate badge
card(c, M, 28 * mm, PW - 2 * M, 22 * mm, CARD_BG)
card(c, M + 4 * mm, 31 * mm, 50 * mm, 16 * mm, ACCENT_GREEN)
txt(c, M + 8 * mm, 35 * mm, "GATE: PASS", size=16, color=WHITE_HEX, bold=True)
txt(c, M + 8 * mm, 43 * mm, "89 / 2 N/A / 0 fail", size=9, color=WHITE_HEX)
txt(c, M + 60 * mm, 36 * mm, "Full system description complete.", size=11, color=WHITE_HEX)
txt(c, M + 60 * mm, 44 * mm, "Ready to proceed to PRD.", size=11, color=ACCENT_GREEN, bold=True)

# Power summary
card(c, M, 54 * mm, PW - 2 * M, 22 * mm, CARD_BG)
txt(c, M + 4 * mm, 58 * mm, "POWER", size=10, color=ACCENT_BLUE, bold=True)
txt(c, M + 4 * mm, 66 * mm, "4.9 mA playing  |  19 uA idle  |  150 mAh battery",
    size=11, color=WHITE_HEX)
txt(c, M + 4 * mm, 73 * mm, "30 hrs continuous  |  10 days typical use  |  45 min charge",
    size=10, color=LIGHT_GRAY)

# Key specs
card(c, M, 80 * mm, PW - 2 * M, 22 * mm, CARD_BG_ALT)
txt(c, M + 4 * mm, 84 * mm, "KEY SPECS", size=10, color=ACCENT_PURPLE, bold=True)
txt(c, M + 4 * mm, 92 * mm, "35x25x10mm pod  |  ~25g total  |  IPX4  |  22mm standard band",
    size=10, color=WHITE_HEX)
txt(c, M + 4 * mm, 99 * mm, "BLE 5.0  |  OTA (signed)  |  USB-C charge  |  $45-60 retail",
    size=10, color=LIGHT_GRAY)

# Open items
bar(c, M, 108 * mm, PW - 2 * M, 1 * mm, ACCENT_ORANGE)
txt(c, M, 112 * mm, "8 OPEN ITEMS", size=11, color=ACCENT_ORANGE, bold=True)

open_items = [
    ("M2", "Haptic perceptibility test with guitarists, pianists, strings"),
    ("M2", "DRV2605L waveform sequence timing precision"),
    ("M2", "Dual-pulse downbeat perception at 120/160/200 BPM"),
    ("M3", "Tap gesture false-trigger rate across instruments"),
    ("M3", "Wristband comfort during 2+ hour sessions"),
    ("M4", "IPX4 seal durability (USB-C flap, 500+ cycles)"),
    ("M1", "BOM volume validation (5k min for $15 target)"),
    ("M3", "Band attachment strength during playing"),
]

for i, (milestone, desc) in enumerate(open_items):
    y_top = 118 * mm + i * 12 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 10 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 3 * mm, milestone, size=9, color=ACCENT_ORANGE, bold=True)
    txt(c, M + 18 * mm, y_top + 3 * mm, desc, size=9, color=SOFT_WHITE)

# CTA
card(c, M, 218 * mm, PW - 2 * M, 16 * mm, CARD_BG)
txt_wrap(c, M + 4 * mm, 222 * mm,
         "Next step: prototype the haptic pod on an nRF52 devkit "
         "with a DRV2605L breakout board. Put it on wrists. Play music. "
         "Validate the three hardest problems before committing to PCB.",
         size=10, color=WHITE_HEX, max_w=PW - 2 * M - 8 * mm, line_h=13)

footer(c, 8, TOTAL_PAGES)
c.showPage()

# ================================================================
c.save()
print(f"Saved {TOTAL_PAGES}-page carousel PDF to {output}")
print(f"Page size: {PW/mm:.0f} x {PH/mm:.0f} mm (4:5 ratio)")

# ================================================================
# PPTX GENERATION
# ================================================================
# 4:5 portrait slide: 7.5 x 9.375 inches (same ratio as PDF)
SLIDE_W = Inches(7.5)
SLIDE_H = Inches(9.375)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]

# -- PPTX helper colors --
C_DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
C_CARD_BG = RGBColor(0x22, 0x22, 0x3A)
C_CARD_BG_ALT = RGBColor(0x1E, 0x1E, 0x34)
C_ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
C_ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0x78)
C_ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
C_ACCENT_BLUE = RGBColor(0x00, 0x9B, 0xF5)
C_ACCENT_PURPLE = RGBColor(0xA8, 0x55, 0xF7)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
C_SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)

PM = Inches(0.4)  # margin


def pptx_bg(slide, color=C_DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def pptx_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def pptx_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def pptx_text(slide, left, top, width, height, text, size=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Helvetica"
    p.alignment = align
    return txBox


def pptx_add_para(tf, text, size=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Helvetica"
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    return p


def pptx_footer(slide, page, total):
    pptx_text(slide, SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.3),
              Inches(0.6), Inches(0.25), f"{page}/{total}",
              size=8, color=C_LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ================================================================
# PPTX PAGE 1: Title
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_PURPLE)

pptx_text(slide, PM, Inches(0.5), SLIDE_W - 2 * PM, Inches(0.6),
          "Haptic Metronome", size=38, bold=True)
pptx_text(slide, PM, Inches(1.0), SLIDE_W - 2 * PM, Inches(0.6),
          "Bracelet", size=38, bold=True)

pptx_rect(slide, PM, Inches(1.5), Inches(1.4), Inches(0.04), C_ACCENT_PURPLE)

pptx_text(slide, PM, Inches(1.7), SLIDE_W - 2 * PM, Inches(0.4),
          "Feel the beat. Hear nothing.", size=18, color=C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(2.2), SLIDE_W - 2 * PM, Inches(0.6),
          "A wrist-worn vibrotactile metronome for musicians who need "
          "silent, precise time -- practice, rehearsal, and stage.",
          size=11, color=C_LIGHT_GRAY)

# Image
img_path = os.path.join(_DIR, "cross_section_illustration_haptic_metronome.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, PM, Inches(3.0),
                             width=SLIDE_W - 2 * PM, height=Inches(5.3))

pptx_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), C_CARD_BG)
pptx_text(slide, PM, SLIDE_H - Inches(0.3), SLIDE_W - 2 * PM, Inches(0.2),
          "Product Overview  |  Concept Stage  |  2026", size=9, color=C_LIGHT_GRAY)

pptx_footer(slide, 1, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 2: The Problem
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_RED)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "The Problem", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_RED)

pptx_problems = [
    ("Audible clicks are unusable on stage",
     "Click tracks bleed into microphones, in-ear monitors isolate "
     "from the room, and audible metronomes distract the audience.",
     C_ACCENT_RED),
    ("Phone apps are imprecise",
     "Android audio latency: 5-40 ms. iOS: 5-15 ms. OS interrupts "
     "and background tasks cause timing jitter. Phone vibration motors "
     "are too slow and vague for musical time.",
     C_ACCENT_ORANGE),
    ("Existing solutions compromise",
     "Earpiece click tracks occupy an ear. Visual metronomes require "
     "looking away from the music. There is no silent, precise, "
     "hands-free time reference that leaves both ears and eyes open.",
     C_ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(pptx_problems):
    y = Inches(1.15 + i * 2.15)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(1.9), C_CARD_BG)
    pptx_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.06), color)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.2), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.3),
              title, size=16, color=color, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.65), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.0),
              desc, size=11, color=C_LIGHT_GRAY)

# Target users
pptx_rounded_rect(slide, PM, Inches(7.7), SLIDE_W - 2 * PM, Inches(1.2), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(7.8), SLIDE_W - 2 * PM, Inches(0.2),
          "TARGET USERS", size=10, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.1), SLIDE_W - 2 * PM, Inches(0.25),
          "Guitarists  |  Pianists  |  String players", size=12, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.45), SLIDE_W - 2 * PM, Inches(0.25),
          "Wind & brass  |  Vocalists  |  Conductors", size=12, bold=True)

pptx_footer(slide, 2, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 3: How It Works
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "How It Works", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_GREEN)

pptx_steps = [
    ("1. Set", "Open app, set BPM and time signature",
     "App pushes config over BLE to the bracelet. "
     "Save presets and setlists for quick recall.",
     C_ACCENT_BLUE),
    ("2. Tap", "Tap the bracelet to start",
     "Accelerometer detects a fingertip tap on the pod. "
     "Haptic pulses begin immediately. LED flashes on beat "
     "for the first 4 bars as a sanity check, then goes dark.",
     C_ACCENT_GREEN),
    ("3. Feel", "Feel the beat on your wrist",
     "Single pulse = normal beat (~15 ms crisp tap). "
     "Dual-pulse = downbeat (\"da-dum\", two taps 40 ms apart). "
     "You always know where beat 1 is.",
     C_ACCENT_PURPLE),
    ("4. Stop", "Double-tap to stop",
     "Accelerometer detects the double-tap gesture. "
     "Pulses stop. Device goes to low-power idle, "
     "ready for the next tap.",
     C_ACCENT_ORANGE),
]

for i, (title, subtitle, desc, color) in enumerate(pptx_steps):
    y = Inches(1.1 + i * 1.95)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(1.7), C_CARD_BG)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.1), Inches(3), Inches(0.35),
              title, size=18, color=color, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.45), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.25),
              subtitle, size=11, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.8), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.8),
              desc, size=10, color=C_LIGHT_GRAY)

pptx_footer(slide, 3, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 4: Architecture
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_BLUE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Architecture", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_BLUE)

# Signal chain
pptx_rounded_rect(slide, PM, Inches(1.1), SLIDE_W - 2 * PM, Inches(1.1), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(1.15), SLIDE_W - 2 * PM, Inches(0.2),
          "Beat Signal Chain", size=11, color=C_ACCENT_BLUE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(1.45), SLIDE_W - 2 * PM, Inches(0.25),
          "Timer ISR  -->  I2C cmd  -->  DRV2605L  -->  LRA motor",
          size=12, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(1.8), SLIDE_W - 2 * PM, Inches(0.2),
          "< 100 us jitter      waveform select     shaped drive     crisp tap on wrist",
          size=8, color=C_LIGHT_GRAY)

# Gesture chain
pptx_rounded_rect(slide, PM, Inches(2.35), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(2.4), SLIDE_W - 2 * PM, Inches(0.2),
          "Gesture Chain", size=11, color=C_ACCENT_GREEN, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(2.7), SLIDE_W - 2 * PM, Inches(0.25),
          "Tap on pod  -->  LIS2DH12 interrupt  -->  FW validation  -->  Start/Stop",
          size=10, bold=True)

pptx_rect(slide, PM, Inches(3.4), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_BLUE)
pptx_text(slide, PM, Inches(3.5), SLIDE_W - 2 * PM, Inches(0.2),
          "SUBSYSTEMS", size=11, color=C_ACCENT_BLUE, bold=True)

pptx_subsystems = [
    ("MCU (nRF52832)", "Timing engine, BLE 5.0, gesture detection, haptic control", C_ACCENT_GREEN),
    ("Haptic (DRV2605L + LRA)", "Crisp vibrotactile pulses, auto-resonance tracking", C_ACCENT_PURPLE),
    ("Accelerometer (LIS2DH12)", "Hardware tap/double-tap detection, Z-axis validation", C_ACCENT_GREEN),
    ("BLE 5.0 (GATT)", "App config, presets, battery status, OTA updates", C_ACCENT_BLUE),
    ("Power (150mAh LiPo)", "USB-C charging, 30+ hrs continuous play", C_ACCENT_ORANGE),
    ("LED (1x green)", "Sanity check: first 4 bars only, then dark", C_ACCENT_ORANGE),
]

for i, (name, desc, color) in enumerate(pptx_subsystems):
    y = Inches(3.8 + i * 0.65)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.55), bg_c)
    pptx_rect(slide, PM, y, Inches(0.06), Inches(0.55), color)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.03), SLIDE_W - 2 * PM, Inches(0.2),
              name, size=11, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.28), SLIDE_W - 2 * PM, Inches(0.2),
              desc, size=8, color=C_LIGHT_GRAY)

# Thin actuator, thick app
pptx_rounded_rect(slide, PM, Inches(7.75), SLIDE_W - 2 * PM, Inches(1.1), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.15), Inches(7.8), SLIDE_W - 2 * PM, Inches(0.25),
          "Thin actuator, thick app", size=13, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.15), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.6),
          "The bracelet handles timing and haptics. The app handles everything "
          "else: BPM, meter, accent patterns, presets, setlists, practice logs. "
          "No physical controls for config -- tap/double-tap for start/stop only.",
          size=9, color=C_LIGHT_GRAY)

pptx_footer(slide, 4, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 5: Downbeat Differentiation
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_PURPLE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Feeling the Downbeat", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_PURPLE)

pptx_text(slide, PM, Inches(1.05), SLIDE_W - 2 * PM, Inches(0.4),
          "How do you know where beat 1 is when you can only feel vibration?",
          size=13, color=C_LIGHT_GRAY)

# Normal beat
pptx_rounded_rect(slide, PM, Inches(1.6), SLIDE_W - 2 * PM, Inches(1.8), C_CARD_BG)
pptx_rect(slide, PM, Inches(1.6), SLIDE_W - 2 * PM, Inches(0.06), C_ACCENT_GREEN)
pptx_text(slide, PM + Inches(0.2), Inches(1.75), SLIDE_W - 2 * PM, Inches(0.3),
          "Normal Beat (2, 3, 4...)", size=16, color=C_ACCENT_GREEN, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(2.2), SLIDE_W - 2 * PM, Inches(0.25),
          "Single pulse: one crisp tap", size=13)
pptx_text(slide, PM + Inches(0.2), Inches(2.6), SLIDE_W - 2 * PM, Inches(0.2),
          "|--15ms--|", size=11, color=C_LIGHT_GRAY)
pptx_rect(slide, PM + Inches(0.2), Inches(2.95), Inches(0.6), Inches(0.12), C_ACCENT_GREEN)

# Downbeat
pptx_rounded_rect(slide, PM, Inches(3.6), SLIDE_W - 2 * PM, Inches(2.15), C_CARD_BG)
pptx_rect(slide, PM, Inches(3.6), SLIDE_W - 2 * PM, Inches(0.06), C_ACCENT_PURPLE)
pptx_text(slide, PM + Inches(0.2), Inches(3.75), SLIDE_W - 2 * PM, Inches(0.3),
          "Downbeat (Beat 1)", size=16, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(4.2), SLIDE_W - 2 * PM, Inches(0.25),
          "Dual-pulse: \"da-dum\"", size=13)
pptx_text(slide, PM + Inches(0.2), Inches(4.55), SLIDE_W - 2 * PM, Inches(0.25),
          "Two taps, 40 ms apart", size=13)
pptx_text(slide, PM + Inches(0.2), Inches(4.95), SLIDE_W - 2 * PM, Inches(0.2),
          "|--15ms--|----40ms----|--15ms--|", size=11, color=C_LIGHT_GRAY)
pptx_rect(slide, PM + Inches(0.2), Inches(5.3), Inches(0.6), Inches(0.12), C_ACCENT_PURPLE)
pptx_rect(slide, PM + Inches(2.35), Inches(5.3), Inches(0.6), Inches(0.12), C_ACCENT_PURPLE)

# Rationale
pptx_rounded_rect(slide, PM, Inches(6.1), SLIDE_W - 2 * PM, Inches(2.75), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.2), Inches(6.2), SLIDE_W - 2 * PM, Inches(0.25),
          "Why pattern, not intensity?", size=13, color=C_ACCENT_ORANGE, bold=True)

txBox = pptx_text(slide, PM + Inches(0.2), Inches(6.6), SLIDE_W - 2 * PM - Inches(0.4), Inches(2.0),
                  "Amplitude is ambiguous: \"Slightly stronger\" vs. \"normal\" is hard to "
                  "distinguish when your arm is moving.", size=9, color=C_LIGHT_GRAY)
tf = txBox.text_frame
pptx_add_para(tf, "Pattern is distinct: The brain detects a double-tap as a different event, "
              "not just a louder version of the same event.",
              size=9, color=C_LIGHT_GRAY, space_before=6)
pptx_add_para(tf, "Works at tempo: At 200 BPM (300 ms between beats), the 70 ms dual-pulse "
              "uses 23% of the interval -- tight but clear.",
              size=9, color=C_LIGHT_GRAY, space_before=6)

pptx_footer(slide, 5, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 6: Constraints & BOM
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_ORANGE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Constraints & BOM", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_ORANGE)

pptx_constraints = [
    ("Timing jitter", "< 100 us", "Hardware timer ISR, not software loop"),
    ("Haptic rise time", "< 10 ms", "LRA mandatory; ERM too slow above 160 BPM"),
    ("Battery life", "> 8 hrs continuous", "150 mAh LiPo, ~4.9 mA avg during play"),
    ("Pod size", "35 x 25 x 10 mm", "Fits comfortably on wrist during playing"),
    ("Pod weight", "< 20g", "Lighter than most watches"),
    ("Sweat resistance", "IPX4", "Silicone overmold, sealed USB-C flap"),
]

for i, (name, value, note) in enumerate(pptx_constraints):
    y = Inches(1.1 + i * 0.65)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.55), bg_c)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.05), Inches(3), Inches(0.2),
              name, size=11, color=C_ACCENT_ORANGE, bold=True)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.3), Inches(5), Inches(0.2),
              note, size=8, color=C_LIGHT_GRAY)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.05), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.2),
              value, size=11, bold=True, align=PP_ALIGN.RIGHT)

# BOM
pptx_rect(slide, PM, Inches(5.15), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_GREEN)
pptx_text(slide, PM, Inches(5.25), SLIDE_W - 2 * PM, Inches(0.2),
          "BOM ESTIMATE (1k units)", size=11, color=C_ACCENT_GREEN, bold=True)

pptx_bom = [
    ("nRF52832 MCU", "$2.80"),
    ("DRV2605L haptic driver", "$1.50"),
    ("LRA motor (8mm)", "$1.50"),
    ("LIS2DH12 accelerometer", "$0.80"),
    ("LiPo 150mAh + protection", "$1.60"),
    ("Crystals + LED + passives", "$0.70"),
    ("USB-C + charge IC + LDO", "$0.95"),
    ("PCB (28x18mm, 4-layer)", "$0.80"),
    ("Pod enclosure (silicone overmold)", "$2.50"),
    ("Silicone wristband + lugs", "$1.20"),
    ("Packaging + USB-C cable", "$1.20"),
    ("Assembly + test", "$2.50"),
]

for i, (item, cost) in enumerate(pptx_bom):
    y = Inches(5.55 + i * 0.27)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rect(slide, PM + Inches(0.08), y, SLIDE_W - 2 * PM - Inches(0.16), Inches(0.23), bg_c)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.02), Inches(4), Inches(0.18),
              item, size=8, color=C_SOFT_WHITE)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.02), SLIDE_W - 2 * PM - Inches(0.6), Inches(0.18),
              cost, size=8, bold=True, align=PP_ALIGN.RIGHT)

# Total
pptx_rect(slide, PM + Inches(0.08), Inches(8.82), SLIDE_W - 2 * PM - Inches(0.16), Inches(0.3), C_ACCENT_ORANGE)
pptx_text(slide, PM + Inches(0.2), Inches(8.84), Inches(2), Inches(0.25),
          "Total COGS", size=10, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(8.84), SLIDE_W - 2 * PM - Inches(0.6), Inches(0.25),
          "~$18.05  (5k: ~$14.50)", size=10, bold=True, align=PP_ALIGN.RIGHT)

pptx_footer(slide, 6, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 7: Hardest Problems
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_RED)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Hardest Problems", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_RED)

pptx_hard = [
    ("Haptic perceptibility during playing",
     "Can the LRA pulse (~1.5g) be felt on the wrist while "
     "strumming guitar, bowing violin, or playing piano expressively? "
     "Arm motion creates competing vibrations. Tight skin contact "
     "and sharp LRA rise time are the primary mitigations. "
     "Must prototype and test with real musicians."),
    ("Tap gesture vs. playing motion",
     "The accelerometer must reliably detect a deliberate fingertip "
     "tap on the pod while rejecting strumming, bowing, and arm "
     "movement. Two-stage detection: hardware threshold trigger, then "
     "firmware Z-axis dominance validation. Target: < 1% false "
     "trigger, < 5% missed tap."),
    ("Downbeat perception at tempo",
     "Can musicians distinguish the dual-pulse downbeat from a "
     "single pulse at 120, 160, and 200 BPM? At 200 BPM the "
     "70 ms dual-pulse occupies 23% of the beat interval. The 40 ms "
     "gap must be above tactile temporal resolution (~5-10 ms). "
     "Needs testing with 5+ musicians across instruments."),
]

for i, (title, desc) in enumerate(pptx_hard):
    y = Inches(1.15 + i * 2.4)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(2.15), C_CARD_BG)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.15), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.25),
              f"{i+1}. {title}", size=13, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.55), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.4),
              desc, size=10, color=C_LIGHT_GRAY)

# Bottom note
pptx_rounded_rect(slide, PM, Inches(8.4), SLIDE_W - 2 * PM, Inches(0.6), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(8.45), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.5),
          "All three require prototype validation with real musicians. "
          "No amount of simulation replaces putting this on a wrist.",
          size=11, color=C_ACCENT_ORANGE)

pptx_footer(slide, 7, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 8: Gate Result & Next
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Gate Result & Next", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_GREEN)

# Gate badge
pptx_rounded_rect(slide, PM, Inches(1.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_rect(slide, PM + Inches(0.15), Inches(1.2), Inches(2.0), Inches(0.6), C_ACCENT_GREEN)
pptx_text(slide, PM + Inches(0.3), Inches(1.22), Inches(1.8), Inches(0.3),
          "GATE: PASS", size=16, bold=True)
pptx_text(slide, PM + Inches(0.3), Inches(1.52), Inches(1.8), Inches(0.2),
          "89 / 2 N/A / 0 fail", size=9)
pptx_text(slide, PM + Inches(2.4), Inches(1.28), Inches(4), Inches(0.2),
          "Full system description complete.", size=11)
pptx_text(slide, PM + Inches(2.4), Inches(1.55), Inches(4), Inches(0.2),
          "Ready to proceed to PRD.", size=11, color=C_ACCENT_GREEN, bold=True)

# Power summary
pptx_rounded_rect(slide, PM, Inches(2.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(2.15), SLIDE_W - 2 * PM, Inches(0.2),
          "POWER", size=10, color=C_ACCENT_BLUE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(2.4), SLIDE_W - 2 * PM, Inches(0.2),
          "4.9 mA playing  |  19 uA idle  |  150 mAh battery", size=11)
pptx_text(slide, PM + Inches(0.15), Inches(2.65), SLIDE_W - 2 * PM, Inches(0.2),
          "30 hrs continuous  |  10 days typical use  |  45 min charge",
          size=10, color=C_LIGHT_GRAY)

# Key specs
pptx_rounded_rect(slide, PM, Inches(3.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.15), Inches(3.15), SLIDE_W - 2 * PM, Inches(0.2),
          "KEY SPECS", size=10, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(3.4), SLIDE_W - 2 * PM, Inches(0.2),
          "35x25x10mm pod  |  ~25g total  |  IPX4  |  22mm standard band", size=10)
pptx_text(slide, PM + Inches(0.15), Inches(3.65), SLIDE_W - 2 * PM, Inches(0.2),
          "BLE 5.0  |  OTA (signed)  |  USB-C charge  |  $45-60 retail",
          size=10, color=C_LIGHT_GRAY)

# Open items
pptx_rect(slide, PM, Inches(4.2), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_ORANGE)
pptx_text(slide, PM, Inches(4.3), SLIDE_W - 2 * PM, Inches(0.2),
          "8 OPEN ITEMS", size=11, color=C_ACCENT_ORANGE, bold=True)

pptx_open_items = [
    ("M2", "Haptic perceptibility test with guitarists, pianists, strings"),
    ("M2", "DRV2605L waveform sequence timing precision"),
    ("M2", "Dual-pulse downbeat perception at 120/160/200 BPM"),
    ("M3", "Tap gesture false-trigger rate across instruments"),
    ("M3", "Wristband comfort during 2+ hour sessions"),
    ("M4", "IPX4 seal durability (USB-C flap, 500+ cycles)"),
    ("M1", "BOM volume validation (5k min for $15 target)"),
    ("M3", "Band attachment strength during playing"),
]

for i, (milestone, desc) in enumerate(pptx_open_items):
    y = Inches(4.6 + i * 0.47)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.38), bg_c)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.07), Inches(0.5), Inches(0.2),
              milestone, size=9, color=C_ACCENT_ORANGE, bold=True)
    pptx_text(slide, PM + Inches(0.7), y + Inches(0.07), SLIDE_W - 2 * PM - Inches(1.0), Inches(0.2),
              desc, size=9, color=C_SOFT_WHITE)

# CTA
pptx_rounded_rect(slide, PM, Inches(8.45), SLIDE_W - 2 * PM, Inches(0.6), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(8.5), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.5),
          "Next step: prototype the haptic pod on an nRF52 devkit "
          "with a DRV2605L breakout board. Put it on wrists. Play music. "
          "Validate the three hardest problems before committing to PCB.",
          size=10)

pptx_footer(slide, 8, TOTAL_PAGES)

# Save PPTX
pptx_output = os.path.join(_DIR, "Haptic_Metronome_Bracelet_Carousel.pptx")
prs.save(pptx_output)
print(f"Saved {TOTAL_PAGES}-slide carousel PPTX to {pptx_output}")
