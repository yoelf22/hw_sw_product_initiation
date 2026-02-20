#!/usr/bin/env python3
"""Build Chair Balancing Act product overview deck."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

_DIR = os.path.dirname(os.path.abspath(__file__))

# -- Theme colors --
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0x78)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_BLUE = RGBColor(0x00, 0x9B, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)
CARD_BG = RGBColor(0x22, 0x22, 0x3A)
CARD_BG_ALT = RGBColor(0x1E, 0x1E, 0x34)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return p


def add_p(tf, text, size=18, color=WHITE, bold=False, before=Pt(6), after=Pt(2), align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    p.space_before = before
    p.space_after = after
    return p


def add_bullet(tf, text, size=16, color=WHITE, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.level = 0
    p.space_before = Pt(4)
    p.space_after = Pt(2)
    return p


def accent_bar(slide, left, top, height=Inches(0.8), color=ACCENT_ORANGE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def circle_num(slide, x, y, num, color):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.4), Inches(0.4))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    ctf = c.text_frame
    ctf.word_wrap = False
    set_text(ctf, str(num), size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_ORANGE)

t = tb(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5))
set_text(t.text_frame, "Chair Balancing Act", size=56, color=WHITE, bold=True)
add_p(t.text_frame, "A tilt-sensing audio feedback device for fixed-leg chairs",
      size=26, color=ACCENT_GREEN, before=Pt(12))

t2 = tb(slide, Inches(1), Inches(3.8), Inches(9), Inches(1.2))
tf = t2.text_frame
set_text(tf, "All legs down = silence.  Any leg lifts = rising tone.  More tilt = louder.", size=22, color=LIGHT_GRAY)
add_p(tf, "Part safety reminder, part office prank, part gift-shop impulse buy.", size=18, color=LIGHT_GRAY, before=Pt(8))

# Bottom
add_shape(slide, Inches(0), Inches(6.5), W, Inches(0.005), RGBColor(0x33, 0x33, 0x55))
t3 = tb(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.6))
set_text(t3.text_frame, "Product Overview  \u2022  Concept Stage  \u2022  2026", size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 2: Product Overview (rendered image)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, color=WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_ORANGE)

t = tb(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7))
set_text(t.text_frame, "Product Overview", size=36, color=DARK_BG, bold=True)

img_path = os.path.join(_DIR, "cBalance.png")
if os.path.exists(img_path):
    img_w = Inches(11.5)
    img_h = Inches(6.2)
    img_left = (W - img_w) // 2
    img_top = Inches(1.0)
    slide.shapes.add_picture(img_path, img_left, img_top, img_w, img_h)

# ============================================================
# SLIDE 3: How It Works + Device
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_GREEN)

t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t.text_frame, "How It Works", size=36, color=WHITE, bold=True)

# Three zones
zones = [
    ("Static Balance", "SILENCE",
     "All legs on the floor.\n\nDevice sleeps.\nNo sound at all.",
     ACCENT_GREEN),
    ("Dynamic Balance", "CONTINUOUS RISING TONE",
     "Any leg lifts off the floor.\n\nAgreeable low tone begins.\nRises smoothly with tilt.\n\nPitch + volume + urgency\nscale proportionally.",
     ACCENT_ORANGE),
    ("Near Fall", "MAXIMUM ALARM",
     "Steep angle, close to\ntipping point.\n\nTone at peak urgency.\nEveryone in the room\nknows what's happening.",
     ACCENT_RED),
]

for i, (title, subtitle, desc, color) in enumerate(zones):
    x = Inches(0.8) + Inches(4.1) * i
    add_shape(slide, x, Inches(1.8), Inches(3.6), Inches(0.06), color)
    add_shape(slide, x, Inches(1.86), Inches(3.6), Inches(3.6), CARD_BG)

    t = tb(slide, x + Inches(0.3), Inches(2.1), Inches(3.0), Inches(0.6))
    set_text(t.text_frame, title, size=24, color=color, bold=True)

    t = tb(slide, x + Inches(0.3), Inches(2.7), Inches(3.0), Inches(0.4))
    set_text(t.text_frame, subtitle, size=12, color=LIGHT_GRAY, bold=True)

    t = tb(slide, x + Inches(0.3), Inches(3.2), Inches(3.0), Inches(2.0))
    set_text(t.text_frame, desc, size=15, color=SOFT_WHITE)

# Bottom — device specs
add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), CARD_BG)
t = tb(slide, Inches(1.0), Inches(5.9), Inches(11), Inches(1.0))
tf = t.text_frame
set_text(tf, "The device:  ", size=16, color=WHITE, bold=True)
add_p(tf, "~32 mm round \u00d7 10 mm  \u2022  ~15g  \u2022  CR2450 coin cell (months of life)  "
      "\u2022  AirTag-sized  \u2022  Adhesive mount under chair seat  \u2022  Fixed-leg chairs only",
      size=15, color=LIGHT_GRAY, before=Pt(4))

# ============================================================
# SLIDE 3: Architecture + Constraints
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t.text_frame, "Architecture & Constraints", size=36, color=WHITE, bold=True)

# Left — subsystems
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.5), ACCENT_BLUE)
t = tb(slide, Inches(1.0), Inches(1.55), Inches(5.4), Inches(0.4))
set_text(t.text_frame, "SUBSYSTEMS", size=14, color=WHITE, bold=True)

subsystems = [
    ("Accelerometer", "Tilt angle on pitch + roll", "I2C + wake interrupt"),
    ("MCU + Firmware", "Filter, tone generation, calibration", "Reads accel, generates PWM tone"),
    ("Class-D Amp + Piezo Speaker", "Amplifies MCU tone output", "PWM input, piezo output"),
    ("CR2450 Coin Cell", "Months of life, user-replaceable", "MCU sleeps in static balance"),
    ("Button", "On/off, mode cycle", "GPIO"),
]

for i, (name, purpose, detail) in enumerate(subsystems):
    y = Inches(2.2) + Inches(0.75) * i
    bg = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    add_shape(slide, Inches(0.8), y, Inches(5.8), Inches(0.65), bg)

    t = tb(slide, Inches(1.0), y + Inches(0.05), Inches(2.2), Inches(0.3))
    set_text(t.text_frame, name, size=14, color=WHITE, bold=True)

    t = tb(slide, Inches(1.0), y + Inches(0.33), Inches(2.2), Inches(0.3))
    set_text(t.text_frame, purpose, size=11, color=LIGHT_GRAY)

    t = tb(slide, Inches(3.5), y + Inches(0.15), Inches(3.0), Inches(0.35))
    set_text(t.text_frame, detail, size=11, color=ACCENT_GREEN)

# Right — constraints
rx = Inches(7.3)
add_shape(slide, rx, Inches(1.5), Inches(5.3), Inches(0.5), ACCENT_ORANGE)
t = tb(slide, rx + Inches(0.2), Inches(1.55), Inches(4.8), Inches(0.4))
set_text(t.text_frame, "HARD CONSTRAINTS", size=14, color=WHITE, bold=True)

constraints = [
    ("BOM cost", "<$8 at 1k units", "Gift/novelty price \u2014 retail $15-25"),
    ("Battery life", ">2 months on CR2450", "Forget-and-use, no charging"),
    ("Latency", "<100 ms tilt-to-sound", "Comedy timing is the product"),
    ("Size", "32 mm \u00d7 10 mm, 15g", "AirTag-sized, hidden under seat"),
]

for i, (name, value, why) in enumerate(constraints):
    y = Inches(2.2) + Inches(0.95) * i
    add_shape(slide, rx, y, Inches(5.3), Inches(0.82), CARD_BG)

    t = tb(slide, rx + Inches(0.2), y + Inches(0.05), Inches(2.0), Inches(0.3))
    set_text(t.text_frame, name, size=14, color=ACCENT_ORANGE, bold=True)

    t = tb(slide, rx + Inches(2.4), y + Inches(0.05), Inches(2.7), Inches(0.3))
    set_text(t.text_frame, value, size=14, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

    t = tb(slide, rx + Inches(0.2), y + Inches(0.4), Inches(4.9), Inches(0.35))
    set_text(t.text_frame, why, size=12, color=LIGHT_GRAY)

# Bottom — data flow
add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8), CARD_BG)
t = tb(slide, Inches(1.0), Inches(6.3), Inches(11), Inches(0.6))
set_text(t.text_frame,
         "Data flow:  Accelerometer --I2C--> MCU (filter + tone gen) --PWM--> Class-D Amp --> Piezo speaker",
         size=15, color=LIGHT_GRAY)

# ============================================================
# SLIDE 4: Hardest Problems + Component Tradeoffs
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_RED)

t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t.text_frame, "Hardest Problems & Component Tradeoffs", size=32, color=WHITE, bold=True)

# Left — three hardest problems
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.5), ACCENT_RED)
t = tb(slide, Inches(1.0), Inches(1.55), Inches(5.4), Inches(0.4))
set_text(t.text_frame, "THREE HARDEST PROBLEMS", size=14, color=WHITE, bold=True)

problems = [
    ("Escalation curve across chair types",
     "School chairs and dining chairs have different fall points and tilt ranges. "
     "The continuous tone must start gently at first leg lift and peak just before the fall point "
     "\u2014 not too early (annoying) or too late (useless). Needs auto-calibration + playtesting."),
    ("Filtering fidgeting from tilting",
     "Sitting down, crossing legs, leaning to grab something \u2014 all produce accelerometer signals. "
     "Continuous model means any sustained tilt produces sound \u2014 the dead-band and time filter are critical."),
    ("Designing the continuous tone",
     "The tilt-to-tone mapping IS the product. Must be agreeable at low tilt and urgent near fall. "
     "Three modes (serious, comedic, stealth) need distinct character. "
     "Needs sound designer + playtesting with kids."),
]

for i, (title, desc) in enumerate(problems):
    y = Inches(2.2) + Inches(1.35) * i
    add_shape(slide, Inches(0.8), y, Inches(5.8), Inches(1.15), CARD_BG)
    circle_num(slide, Inches(1.0), y + Inches(0.15), i + 1, ACCENT_RED)

    t = tb(slide, Inches(1.6), y + Inches(0.08), Inches(4.8), Inches(0.35))
    set_text(t.text_frame, title, size=15, color=WHITE, bold=True)

    t = tb(slide, Inches(1.6), y + Inches(0.42), Inches(4.8), Inches(0.7))
    set_text(t.text_frame, desc, size=11, color=LIGHT_GRAY)

# Right — component tradeoffs
rx = Inches(7.3)
add_shape(slide, rx, Inches(1.5), Inches(5.3), Inches(0.5), ACCENT_GREEN)
t = tb(slide, rx + Inches(0.2), Inches(1.55), Inches(4.8), Inches(0.4))
set_text(t.text_frame, "COMPONENT TRADEOFFS", size=14, color=WHITE, bold=True)

tradeoffs = [
    ("MCU", "Cost + PWM", "ATtiny ($0.50) has HW PWM for tone gen. STM32L0 ($1.20) has DAC for smoother audio. nRF52 ($3) overkill."),
    ("Accelerometer", "FW complexity", "Cheapest has no filtering \u2014 $0.50 more gets on-chip motion detection, saves weeks of FW tuning."),
    ("Audio path", "Simplicity", "MCU PWM + class-D amp ($0.30) replaces audio playback IC. Enables real-time continuous tone, not fixed clips."),
    ("Speaker", "Physical constraint", "Dynamic (4 mm tall, louder) vs. piezo (2.5 mm, fits 10 mm enclosure). Piezo resonance aids escalation."),
    ("Power", "Simplicity", "CR2450 coin cell: no charging circuit, no USB port, months of life, user-replaceable. LiPo only if V2 needs more current."),
]

for i, (component, axis, desc) in enumerate(tradeoffs):
    y = Inches(2.2) + Inches(0.88) * i
    bg = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    add_shape(slide, rx, y, Inches(5.3), Inches(0.76), bg)

    t = tb(slide, rx + Inches(0.2), y + Inches(0.05), Inches(2.2), Inches(0.3))
    set_text(t.text_frame, component, size=13, color=WHITE, bold=True)

    t = tb(slide, rx + Inches(2.6), y + Inches(0.05), Inches(2.5), Inches(0.3))
    set_text(t.text_frame, axis, size=12, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.RIGHT)

    t = tb(slide, rx + Inches(0.2), y + Inches(0.33), Inches(4.9), Inches(0.4))
    set_text(t.text_frame, desc, size=10, color=LIGHT_GRAY)

# Bottom
t = tb(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6))
set_text(t.text_frame,
         "Hardware is simple. The product lives or dies on the escalation curve \u2014 how tilt maps to tone. Playtesting > engineering.",
         size=15, color=LIGHT_GRAY)

# ============================================================
# SLIDE 5: PRD Summary — Requirements at a Glance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t.text_frame, "PRD Summary: 45 Requirements", size=36, color=WHITE, bold=True)

# Top-left — requirement counts by track
add_shape(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(0.5), ACCENT_BLUE)
t = tb(slide, Inches(1.0), Inches(1.55), Inches(3.2), Inches(0.4))
set_text(t.text_frame, "BY TRACK", size=14, color=WHITE, bold=True)

tracks = [
    ("Hardware", "18 reqs", "(8 electrical, 8 mechanical, 4 PCB)", ACCENT_ORANGE),
    ("Firmware", "18 reqs", "(7 core, 4 modes, 4 calibration, 2 power, 1 versioning)", ACCENT_GREEN),
    ("Integration", "9 reqs", "(end-to-end tests spanning HW + FW)", ACCENT_RED),
]

for i, (name, count, detail, color) in enumerate(tracks):
    y = Inches(2.2) + Inches(0.85) * i
    add_shape(slide, Inches(0.8), y, Inches(3.5), Inches(0.72), CARD_BG)
    accent_bar(slide, Inches(0.8), y + Inches(0.1), Inches(0.52), color)

    t = tb(slide, Inches(1.1), y + Inches(0.05), Inches(1.6), Inches(0.3))
    set_text(t.text_frame, name, size=16, color=WHITE, bold=True)

    t = tb(slide, Inches(2.8), y + Inches(0.05), Inches(1.3), Inches(0.3))
    set_text(t.text_frame, count, size=16, color=color, bold=True, align=PP_ALIGN.RIGHT)

    t = tb(slide, Inches(1.1), y + Inches(0.38), Inches(3.0), Inches(0.3))
    set_text(t.text_frame, detail, size=11, color=LIGHT_GRAY)

# Top-right — priority split
rx = Inches(5.0)
add_shape(slide, rx, Inches(1.5), Inches(3.5), Inches(0.5), ACCENT_ORANGE)
t = tb(slide, rx + Inches(0.2), Inches(1.55), Inches(3.2), Inches(0.4))
set_text(t.text_frame, "BY PRIORITY", size=14, color=WHITE, bold=True)

add_shape(slide, rx, Inches(2.2), Inches(3.5), Inches(1.27), CARD_BG)
t = tb(slide, rx + Inches(0.2), Inches(2.3), Inches(3.0), Inches(0.4))
set_text(t.text_frame, "40 Must-have", size=20, color=ACCENT_ORANGE, bold=True)
t = tb(slide, rx + Inches(0.2), Inches(2.8), Inches(3.0), Inches(0.4))
set_text(t.text_frame, "5 Should-have", size=20, color=ACCENT_GREEN, bold=True)
t = tb(slide, rx + Inches(0.2), Inches(3.2), Inches(3.0), Inches(0.3))
set_text(t.text_frame, "0 Nice-to-have (V2 only)", size=14, color=LIGHT_GRAY)

# Far-right — gate result
gx = Inches(9.2)
add_shape(slide, gx, Inches(1.5), Inches(3.5), Inches(0.5), ACCENT_GREEN)
t = tb(slide, gx + Inches(0.2), Inches(1.55), Inches(3.2), Inches(0.4))
set_text(t.text_frame, "GATE RESULT", size=14, color=WHITE, bold=True)

add_shape(slide, gx, Inches(2.2), Inches(3.5), Inches(1.27), CARD_BG)
t = tb(slide, gx + Inches(0.2), Inches(2.35), Inches(3.0), Inches(0.4))
set_text(t.text_frame, "PASS", size=28, color=ACCENT_GREEN, bold=True)
t = tb(slide, gx + Inches(0.2), Inches(2.9), Inches(3.0), Inches(0.5))
set_text(t.text_frame, "66 pass / 23 N/A / 1 minor gap\n(FW versioning scheme)", size=12, color=LIGHT_GRAY)

# Bottom — key requirements highlights
add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.9), Inches(0.5), ACCENT_GREEN)
t = tb(slide, Inches(1.0), Inches(4.55), Inches(11.5), Inches(0.4))
set_text(t.text_frame, "KEY REQUIREMENTS", size=14, color=WHITE, bold=True)

key_reqs = [
    ("FW-C-05", "Continuous tone: pitch + volume scale proportionally with tilt"),
    ("FW-C-06", "Real-time tracking: angle change to tone change < 100ms"),
    ("INT-02", "Zero false triggers on transient motion (sitting down, bumping)"),
    ("FW-PM-04", "Battery life >= 6 months at 20 tilt events/day"),
    ("HW-M-01", "Enclosure max 35mm dia x 10mm height"),
    ("HW-E-06", "Total sleep current <= 2 uA (static balance)"),
]

for i, (req_id, desc) in enumerate(key_reqs):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + Inches(6.1) * col
    y = Inches(5.2) + Inches(0.55) * row
    bg = CARD_BG if row % 2 == 0 else CARD_BG_ALT
    add_shape(slide, x, y, Inches(5.8), Inches(0.45), bg)

    t = tb(slide, x + Inches(0.15), y + Inches(0.06), Inches(0.9), Inches(0.3))
    set_text(t.text_frame, req_id, size=11, color=ACCENT_BLUE, bold=True)

    t = tb(slide, x + Inches(1.1), y + Inches(0.06), Inches(4.5), Inches(0.3))
    set_text(t.text_frame, desc, size=12, color=SOFT_WHITE)

# ============================================================
# SLIDE 6: Open Items & V2 Horizon
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_ORANGE)

t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t.text_frame, "Open Items & V2 Horizon", size=36, color=WHITE, bold=True)

# Left — 7 open items
add_shape(slide, Inches(0.8), Inches(1.5), Inches(6.5), Inches(0.5), ACCENT_ORANGE)
t = tb(slide, Inches(1.0), Inches(1.55), Inches(6.0), Inches(0.4))
set_text(t.text_frame, "7 OPEN ITEMS (resolve before production)", size=14, color=WHITE, bold=True)

open_items = [
    ("M2", "Escalation curve tuning across 5+ chair types"),
    ("M2", "Piezo volume validation in real environments"),
    ("M2", "False trigger testing with real users (fidgety kids)"),
    ("M2", "PWM audio quality -- pleasant tone at low tilt?"),
    ("M2", "CPSIA / EN 71 acoustic safety for children's products"),
    ("M3", "Tone design for 3 modes (sound designer + FW)"),
    ("M3", "Adhesive durability testing"),
]

for i, (milestone, desc) in enumerate(open_items):
    y = Inches(2.2) + Inches(0.62) * i
    bg = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    add_shape(slide, Inches(0.8), y, Inches(6.5), Inches(0.52), bg)

    t = tb(slide, Inches(1.0), y + Inches(0.1), Inches(0.6), Inches(0.3))
    set_text(t.text_frame, milestone, size=12, color=ACCENT_ORANGE, bold=True)

    t = tb(slide, Inches(1.7), y + Inches(0.1), Inches(5.4), Inches(0.3))
    set_text(t.text_frame, desc, size=13, color=SOFT_WHITE)

# Right — V2 features
rx = Inches(8.0)
add_shape(slide, rx, Inches(1.5), Inches(4.7), Inches(0.5), ACCENT_BLUE)
t = tb(slide, rx + Inches(0.2), Inches(1.55), Inches(4.3), Inches(0.4))
set_text(t.text_frame, "V2 (if market validates)", size=14, color=WHITE, bold=True)

v2_features = [
    ("Custom tone profiles via USB", "Board respin: USB-C connector,\nupload escalation curve parameters"),
    ("BLE + companion app", "MCU change: ATtiny -> nRF52.\nFull board respin + app dev"),
    ("Adjustable sensitivity", "Requires BLE.\nDepends on V2 BLE decision"),
    ("Rechargeable battery", "Adds charging IC, USB-C, Li-Po.\nEnclosure redesign"),
]

for i, (feature, impact) in enumerate(v2_features):
    y = Inches(2.2) + Inches(1.1) * i
    add_shape(slide, rx, y, Inches(4.7), Inches(0.95), CARD_BG)

    t = tb(slide, rx + Inches(0.2), y + Inches(0.05), Inches(4.3), Inches(0.3))
    set_text(t.text_frame, feature, size=14, color=ACCENT_BLUE, bold=True)

    t = tb(slide, rx + Inches(0.2), y + Inches(0.37), Inches(4.3), Inches(0.5))
    set_text(t.text_frame, impact, size=11, color=LIGHT_GRAY)

# Bottom
add_shape(slide, Inches(0.8), Inches(6.6), Inches(11.9), Inches(0.5), CARD_BG)
t = tb(slide, Inches(1.0), Inches(6.65), Inches(11.5), Inches(0.4))
set_text(t.text_frame,
         "V2 decisions should wait for 6+ months of V1 sales data and customer feedback.",
         size=14, color=LIGHT_GRAY)

# ============================================================
# Save
# ============================================================
output_path = os.path.join(_DIR, "Chair_Balancing_Act_Deck.pptx")
prs.save(output_path)
print(f"Saved to {output_path}")
