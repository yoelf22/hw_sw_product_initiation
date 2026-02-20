#!/usr/bin/env python3
"""Build AirSense High-Level System Design deck — executive summary format.

5 slides mirroring the single-page hw_sw_high_level document:
  1. System Overview — concept art + product description
  2. System Architecture — block diagram + hardware cross-sections
  3. Subsystems + Key Interfaces + Constraints
  4. Fundamental HW Problems + Component Choice Architecture
  5. Three Hardest Problems + Open Calls
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

_DIR = os.path.dirname(os.path.abspath(__file__))

# -- Theme --
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
TEAL = RGBColor(0x00, 0xBF, 0xA5)
BLUE = RGBColor(0x00, 0x9B, 0xF5)
PURPLE = RGBColor(0x9B, 0x6D, 0xFF)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED = RGBColor(0xFF, 0x45, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xBB, 0xBB, 0xCC)
SOFT = RGBColor(0xF0, 0xF0, 0xF5)
CARD = RGBColor(0x18, 0x22, 0x3A)
CARD2 = RGBColor(0x14, 0x1E, 0x34)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = DARK_BG

def box(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    return s

def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)

def txt(slide, l, t, w, h, text, sz=14, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    tf = tb(slide, l, t, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = "Calibri"; p.alignment = a
    return tf

def add_p(tf, text, sz=14, c=WHITE, b=False, before=Pt(4), after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = "Calibri"
    p.space_before = before; p.space_after = after
    return p

def accent(slide, l, t, h, c):
    box(slide, l, t, Inches(0.07), h, c)

def num_circle(slide, x, y, n, c):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.36), Inches(0.36))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(n); p.font.size = Pt(14)
    p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

def strip(slide, l, t, w, c):
    box(slide, l, t, w, Inches(0.05), c)

def section_hdr(slide, l, t, w, label, c):
    box(slide, l, t, w, Inches(0.42), c)
    txt(slide, l + Inches(0.15), t + Inches(0.05), w - Inches(0.3), Inches(0.3),
        label, 12, WHITE, True)


# ================================================================
# SLIDE 1 — System Overview (concept art + description)
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
strip(sl, 0, 0, W, TEAL)

# Title
txt(sl, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
    "AirSense  --  High-Level System Design", 32, WHITE, True)
txt(sl, Inches(0.8), Inches(0.85), Inches(11), Inches(0.3),
    "Single-page executive summary  |  Draft  |  2026-02-20", 12, GRAY)

# System overview concept art — left half
overview_img = os.path.join(_DIR, "System_Overview.png")
if os.path.exists(overview_img):
    sl.shapes.add_picture(overview_img, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5))

# "What It Is" card — right half
box(sl, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.4), CARD)
accent(sl, Inches(6.8), Inches(1.5), Inches(2.4), TEAL)
txt(sl, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.3),
    "WHAT IT IS", 12, TEAL, True)
txt(sl, Inches(7.1), Inches(2.0), Inches(5.2), Inches(1.7),
    "A wireless indoor environment monitor that tracks CO2, temperature, humidity, and particulate matter "
    "per room in commercial offices. Battery-powered BLE sensor nodes communicate to per-floor Ethernet "
    "gateways, which forward data to a cloud backend serving a web dashboard for facility managers. "
    "Replaces complaint-driven HVAC management with room-level data and historical trends.",
    13, SOFT)

# Key facts — right half below description
box(sl, Inches(6.8), Inches(4.15), Inches(5.7), Inches(2.8), CARD)
accent(sl, Inches(6.8), Inches(4.15), Inches(2.8), BLUE)
txt(sl, Inches(7.1), Inches(4.25), Inches(5.2), Inches(0.3),
    "AT A GLANCE", 12, BLUE, True)

facts = [
    ("Deployment", "10-200 rooms per building, self-installed magnetic mount"),
    ("Three tiers", "Battery sensor nodes -> Ethernet gateways -> cloud dashboard"),
    ("Battery life", ">12 months on 2x AA lithium (mount and forget)"),
    ("Key sensors", "CO2 (SCD41), temp/humidity (SHT40), PM (PMSA003I)"),
    ("Connectivity", "BLE 5.3 advertising -> MQTT/TLS over Ethernet"),
    ("Target cost", "<$35 sensor node, <$40 gateway at 1k units"),
]

for i, (label, value) in enumerate(facts):
    y = Inches(4.65) + Inches(0.35) * i
    bgc = CARD if i % 2 == 0 else CARD2
    txt(sl, Inches(7.1), y, Inches(1.6), Inches(0.25),
        label, 10, TEAL, True)
    txt(sl, Inches(8.8), y, Inches(3.5), Inches(0.25),
        value, 10, SOFT)

# Footer
txt(sl, Inches(0.8), Inches(7.1), Inches(11), Inches(0.3),
    "HIGH-LEVEL DESIGN  |  Not a PRD  |  Details in system_description_smart_sensor_hub.md",
    10, GRAY)


# ================================================================
# SLIDE 2 — Block Diagram + Cross-Section
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
strip(sl, 0, 0, W, BLUE)

txt(sl, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
    "System Architecture", 32, WHITE, True)

# Block diagram — top
img_path = os.path.join(_DIR, "AirSense_Block_Diagram.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.8), Inches(1.2), Inches(8.0), Inches(3.9))

# Sensor node cross-section — right
node_img = os.path.join(_DIR, "Cross-Section — Sensor Node.png")
if os.path.exists(node_img):
    sl.shapes.add_picture(node_img, Inches(9.2), Inches(1.2), Inches(3.5), Inches(3.5))

# Gateway cross-section — bottom right
gw_img = os.path.join(_DIR, "Cross-Section — Gateway.png")
if os.path.exists(gw_img):
    sl.shapes.add_picture(gw_img, Inches(9.2), Inches(4.9), Inches(3.5), Inches(2.2))

# Labels for cross-sections
box(sl, Inches(9.2), Inches(4.6), Inches(3.5), Inches(0.3), CARD)
txt(sl, Inches(9.3), Inches(4.62), Inches(3.3), Inches(0.25),
    "SENSOR NODE", 10, TEAL, True, PP_ALIGN.CENTER)

box(sl, Inches(9.2), Inches(7.1), Inches(3.5), Inches(0.3), CARD)
txt(sl, Inches(9.3), Inches(7.12), Inches(3.3), Inches(0.25),
    "GATEWAY", 10, BLUE, True, PP_ALIGN.CENTER)

# Architecture summary — bottom left
box(sl, Inches(0.8), Inches(5.3), Inches(8.0), Inches(1.8), CARD)
accent(sl, Inches(0.8), Inches(5.3), Inches(1.8), TEAL)
txt(sl, Inches(1.1), Inches(5.4), Inches(7.5), Inches(0.3),
    "THREE-TIER ARCHITECTURE", 12, TEAL, True)

tf = txt(sl, Inches(1.1), Inches(5.8), Inches(7.5), Inches(1.2),
    "Sensor nodes (1 per room) wake every 5 minutes, read sensors, and broadcast a 20-byte BLE "
    "advertisement. Gateways (1 per floor) passively scan BLE and forward aggregated readings to "
    "the cloud via MQTT over Ethernet.", 11, SOFT)
add_p(tf, "Data flows one direction: sensor -> gateway -> cloud -> dashboard. "
    "OTA firmware updates reverse the path: cloud -> gateway -> BLE DFU to sleeping nodes.",
    11, GRAY, before=Pt(8))


# ================================================================
# SLIDE 2 — Subsystems + Interfaces + Constraints
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
strip(sl, 0, 0, W, BLUE)

txt(sl, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
    "Subsystems, Interfaces & Constraints", 32, WHITE, True)

# -- Left: Subsystems --
lx = Inches(0.8)
section_hdr(sl, lx, Inches(1.3), Inches(4.0), "SUBSYSTEMS", TEAL)

subsystems = [
    ("Sensor array", "CO2, T/H, PM per room", "HW", TEAL),
    ("nRF52840 + FW", "Sensor mgr, BLE adv, power, OTA", "FW", BLUE),
    ("Power (TPS62740 + 2xAA)", "6V -> 3.3V, >12 month target", "HW", ORANGE),
    ("BLE Gateway (ESP32-S3)", "Scan, aggregate, MQTT uplink", "HW+FW", BLUE),
    ("Cloud backend", "MQTT, TimescaleDB, REST, alerts", "Cloud", PURPLE),
    ("Dashboard + QR page", "Floor map, single-room view", "SW", PURPLE),
]

for i, (name, purpose, domain, c) in enumerate(subsystems):
    y = Inches(1.85) + Inches(0.52) * i
    bgc = CARD if i % 2 == 0 else CARD2
    box(sl, lx, y, Inches(4.0), Inches(0.44), bgc)
    accent(sl, lx, y + Inches(0.05), Inches(0.34), c)
    txt(sl, lx + Inches(0.18), y + Inches(0.05), Inches(1.9), Inches(0.2),
        name, 10, WHITE, True)
    txt(sl, lx + Inches(0.18), y + Inches(0.24), Inches(2.5), Inches(0.18),
        purpose, 8, GRAY)
    txt(sl, lx + Inches(3.2), y + Inches(0.1), Inches(0.7), Inches(0.2),
        domain, 8, c, True, PP_ALIGN.RIGHT)

# -- Center: Key Interfaces --
cx = Inches(5.2)
section_hdr(sl, cx, Inches(1.3), Inches(3.9), "KEY INTERFACES", BLUE)

interfaces = [
    ("Sensors -> MCU", "I2C (100-400 kHz)", "Raw CO2, T, RH, PM"),
    ("MCU -> Gateway", "BLE 5.3 ext adv", "20-byte payload, one-way"),
    ("Gateway -> Cloud", "MQTT / TLS, Ethernet", "JSON telemetry per device"),
    ("Cloud -> Dashboard", "HTTPS REST", "Room data, alerts, health"),
    ("Cloud -> GW -> Node", "MQTT + BLE DFU", "OTA firmware images"),
]

for i, (route, proto, data) in enumerate(interfaces):
    y = Inches(1.85) + Inches(0.62) * i
    bgc = CARD if i % 2 == 0 else CARD2
    box(sl, cx, y, Inches(3.9), Inches(0.54), bgc)
    txt(sl, cx + Inches(0.12), y + Inches(0.04), Inches(2.0), Inches(0.2),
        route, 10, WHITE, True)
    txt(sl, cx + Inches(0.12), y + Inches(0.22), Inches(1.5), Inches(0.15),
        proto, 8, BLUE, True)
    txt(sl, cx + Inches(1.7), y + Inches(0.22), Inches(2.0), Inches(0.28),
        data, 8, GRAY)

# -- Right: Constraints --
rx = Inches(9.5)
section_hdr(sl, rx, Inches(1.3), Inches(3.5), "HARD CONSTRAINTS", ORANGE)

constraints = [
    ("Battery life", ">12 months on 2xAA", "Drives BLE-not-WiFi decision"),
    ("Node BOM", "<$35 at 1k", "50-200 nodes per building"),
    ("Gateway BOM", "<$40 at 1k", "1 per floor"),
    ("CO2 accuracy", "+/-50 ppm + 5%", "HVAC decisions depend on this"),
    ("Certification", "FCC, CE, IC", "BLE = intentional radiator"),
    ("Density", "50 nodes / gateway", "No packet loss at scale"),
]

for i, (name, value, why) in enumerate(constraints):
    y = Inches(1.85) + Inches(0.52) * i
    bgc = CARD if i % 2 == 0 else CARD2
    box(sl, rx, y, Inches(3.5), Inches(0.44), bgc)
    txt(sl, rx + Inches(0.1), y + Inches(0.03), Inches(1.3), Inches(0.18),
        name, 9, ORANGE, True)
    txt(sl, rx + Inches(1.5), y + Inches(0.03), Inches(1.8), Inches(0.18),
        value, 9, WHITE, True, PP_ALIGN.RIGHT)
    txt(sl, rx + Inches(0.1), y + Inches(0.24), Inches(3.2), Inches(0.18),
        why, 8, GRAY)

# Bottom summary
box(sl, Inches(0.8), Inches(5.2), Inches(12.2), Inches(0.5), CARD)
txt(sl, Inches(1.0), Inches(5.25), Inches(11.8), Inches(0.4),
    "Three tiers: battery BLE nodes (per room) -> mains-powered Ethernet gateways (per floor) -> "
    "cloud backend (centralized).  The 12-month battery target is the constraint that shapes everything.",
    12, GRAY)

# Slide number
txt(sl, Inches(0.8), Inches(7.1), Inches(11), Inches(0.3),
    "6 subsystems  |  5 key interfaces  |  6 hard constraints", 10, GRAY)


# ================================================================
# SLIDE 3 — Fundamental HW Problems + Component Choice Architecture
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
strip(sl, 0, 0, W, RED)

txt(sl, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
    "Fundamental HW Problems & Component Tradeoffs", 32, WHITE, True)

# -- Left: Fundamental HW Problems --
section_hdr(sl, Inches(0.8), Inches(1.3), Inches(5.8), "FUNDAMENTAL HARDWARE PROBLEMS", RED)

hw_problems = [
    ("Powering a PM sensor from AA for 12 months",
     "PMSA003I draws 25 mA during sampling -- 50x more than all other components combined. "
     "Defines the entire power architecture. If unsolvable, the product either drops PM "
     "sensing or drops the 12-month battery target."),
    ("Accurate CO2 with a 5-second sensing window",
     "SCD41 specifies 15s warm-up for full accuracy. A 5s window (to save power) may "
     "degrade accuracy below +/-50 ppm. If 15s is required, the CO2 power budget triples."),
    ("BLE reliability at 50 nodes per gateway",
     "50 nodes advertising on 3 BLE channels creates collision probability. "
     "Missed advertisements mean dashboard gaps that undermine facility manager trust."),
]

for i, (title, desc) in enumerate(hw_problems):
    y = Inches(1.85) + Inches(1.55) * i
    box(sl, Inches(0.8), y, Inches(5.8), Inches(1.35), CARD)
    num_circle(sl, Inches(1.0), y + Inches(0.12), i + 1, RED)
    txt(sl, Inches(1.55), y + Inches(0.1), Inches(4.8), Inches(0.28),
        title, 13, WHITE, True)
    txt(sl, Inches(1.55), y + Inches(0.42), Inches(4.8), Inches(0.85),
        desc, 10, GRAY)

# -- Right: Component Choice Architecture --
section_hdr(sl, Inches(7.0), Inches(1.3), Inches(5.7), "COMPONENT CHOICE ARCHITECTURE", TEAL)

components = [
    ("nRF52840", "Performance",
     "Best BLE sleep ($4.50) vs. ESP32-C3 ($1.50, worse sleep). 12-month battery forces the premium.",
     TEAL),
    ("SCD41 CO2", "Performance",
     "Only photoacoustic CO2 at this size + accuracy. $10 dominates BOM. No cheaper alternative.",
     TEAL),
    ("PMSA003I PM", "Cost vs. power",
     "Adds $8 + destroys battery budget. Removable module (two SKUs) lets customer decide.",
     ORANGE),
    ("TPS62740 reg", "Power efficiency",
     "360 nA quiescent ($1.50) vs. LDO ($0.30) that wastes months of battery in quiescent.",
     ORANGE),
    ("ESP32-S3 (GW)", "Cost",
     "Cheapest BLE + Ethernet path. Mains-powered -- no power tension.",
     BLUE),
]

for i, (name, axis, desc, c) in enumerate(components):
    y = Inches(1.85) + Inches(0.95) * i
    bgc = CARD if i % 2 == 0 else CARD2
    box(sl, Inches(7.0), y, Inches(5.7), Inches(0.82), bgc)
    accent(sl, Inches(7.0), y + Inches(0.08), Inches(0.66), c)

    txt(sl, Inches(7.2), y + Inches(0.06), Inches(2.2), Inches(0.22),
        name, 11, WHITE, True)
    txt(sl, Inches(9.6), y + Inches(0.06), Inches(2.8), Inches(0.22),
        axis, 10, c, True, PP_ALIGN.RIGHT)
    txt(sl, Inches(7.2), y + Inches(0.33), Inches(5.2), Inches(0.42),
        desc, 9, GRAY)

# Bottom
box(sl, Inches(0.8), Inches(6.7), Inches(11.9), Inches(0.45), CARD)
txt(sl, Inches(1.0), Inches(6.75), Inches(11.5), Inches(0.35),
    "The PM sensor is the central tension: it differentiates the product but threatens the battery target. "
    "Two-SKU model (with/without PM) is the likely resolution.",
    11, ORANGE)


# ================================================================
# SLIDE 4 — Three Hardest Problems + Open Calls
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
strip(sl, 0, 0, W, PURPLE)

txt(sl, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
    "Hardest Problems & Open Calls", 32, WHITE, True)

# -- Left: Three Hardest Problems --
section_hdr(sl, Inches(0.8), Inches(1.3), Inches(6.8), "THREE HARDEST PROBLEMS", PURPLE)

hardest = [
    ("PM sensor vs. 12-month battery",
     "PMSA003I at 5-min intervals alone exceeds the power budget. "
     "Options: reduce PM sampling to 30-60 min, make PM a removable "
     "module (two SKUs), or accept shorter battery for PM units.",
     RED),
    ("BLE range in concrete offices",
     "Advertising reach drops to 5-10m through concrete walls. "
     "1 gateway per floor assumes 15m LOS. Dense offices may need "
     "2-3 gateways per floor, increasing deployment cost.",
     ORANGE),
    ("OTA firmware updates over BLE via gateway",
     "Cloud -> gateway (MQTT) -> sensor node (BLE DFU) to a device "
     "that sleeps 99% of the time. Must coordinate wake windows, "
     "handle interrupted transfers, never brick deployed devices.",
     BLUE),
]

for i, (title, desc, c) in enumerate(hardest):
    y = Inches(1.85) + Inches(1.55) * i
    box(sl, Inches(0.8), y, Inches(6.8), Inches(1.35), CARD)
    num_circle(sl, Inches(1.0), y + Inches(0.15), i + 1, c)
    txt(sl, Inches(1.55), y + Inches(0.12), Inches(5.8), Inches(0.3),
        title, 14, WHITE, True)
    txt(sl, Inches(1.55), y + Inches(0.48), Inches(5.8), Inches(0.8),
        desc, 11, GRAY)

# -- Right: Open Calls --
section_hdr(sl, Inches(8.0), Inches(1.3), Inches(4.7), "OPEN CALLS (block detailed design)", ORANGE)

calls = [
    ("PM sampling strategy",
     "Every 5 min (kills battery) vs.\nevery 30-60 min (low resolution) vs.\nremovable module (two SKUs)",
     "Before power arch"),
    ("SCD41 sensing window",
     "5s (power-optimized, accuracy risk)\nvs. 15s (spec-compliant, 3x power)\nvs. duty-cycled (TBD accuracy)",
     "Before FW sensor mgr"),
    ("OTA mechanism",
     "BLE DFU via gateway (complex)\nvs. USB-C on device (simple,\nrequires physical access)",
     "Before FW partition"),
    ("Gateway density",
     "1 per floor (cheaper, range risk)\nvs. 1 per zone (reliable,\nmore hardware)",
     "Before deploy guide"),
]

for i, (decision, options, deadline) in enumerate(calls):
    y = Inches(1.85) + Inches(1.2) * i
    bgc = CARD if i % 2 == 0 else CARD2
    box(sl, Inches(8.0), y, Inches(4.7), Inches(1.05), bgc)
    accent(sl, Inches(8.0), y + Inches(0.1), Inches(0.85), ORANGE)

    txt(sl, Inches(8.2), y + Inches(0.06), Inches(4.3), Inches(0.22),
        decision, 11, WHITE, True)
    txt(sl, Inches(8.2), y + Inches(0.3), Inches(3.0), Inches(0.65),
        options, 9, SOFT)
    txt(sl, Inches(11.0), y + Inches(0.06), Inches(1.5), Inches(0.22),
        deadline, 8, ORANGE, True, PP_ALIGN.RIGHT)

# Bottom
box(sl, Inches(0.8), Inches(6.7), Inches(11.9), Inches(0.45), CARD)
txt(sl, Inches(1.0), Inches(6.75), Inches(11.5), Inches(0.35),
    "These 4 decisions must be resolved before the full system description can proceed. "
    "All are interconnected -- the PM sampling decision cascades into power, BOM, and SKU strategy.",
    11, GRAY)


# ================================================================
output = os.path.join(_DIR, "AirSense_High_Level_Deck.pptx")
prs.save(output)
print(f"Saved to {output}")
