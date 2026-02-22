#!/usr/bin/env python3
"""Build AirSense Indoor Environment Monitor product overview deck."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

_DIR = os.path.dirname(os.path.abspath(__file__))

# -- Theme colors --
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_TEAL = RGBColor(0x00, 0xBF, 0xA5)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_BLUE = RGBColor(0x00, 0x9B, 0xF5)
ACCENT_PURPLE = RGBColor(0x9B, 0x6D, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)
CARD_BG = RGBColor(0x18, 0x22, 0x3A)
CARD_BG_ALT = RGBColor(0x14, 0x1E, 0x34)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return p


def add_p(tf, text, size=18, color=WHITE, bold=False, before=Pt(6), after=Pt(2), align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    p.space_before = before
    p.space_after = after
    return p


def accent_bar(slide, left, top, height=Inches(0.8), color=ACCENT_TEAL):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def circle_num(slide, x, y, num, color):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.4), Inches(0.4))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    ctf = c.text_frame
    ctf.word_wrap = False
    set_text(ctf, str(num), size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_TEAL)

t = tb(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(1.5))
set_text(t.text_frame, "AirSense", size=56, color=WHITE, bold=True)
add_p(t.text_frame, "Indoor Environment Monitor",
      size=26, color=ACCENT_TEAL, before=Pt(12))

t2 = tb(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(2.0))
tf = t2.text_frame
set_text(tf, "Room-level CO2, temperature, humidity, and particulate matter.", size=22, color=LIGHT_GRAY)
add_p(tf, "Battery-powered BLE sensor nodes. Per-floor gateways. Cloud dashboard.", size=18, color=LIGHT_GRAY, before=Pt(8))
add_p(tf, "Replace complaint-driven HVAC management with data.", size=18, color=ACCENT_ORANGE, before=Pt(8))

# System overview concept art — right side
overview_img = os.path.join(_DIR, "System_Overview.png")
if os.path.exists(overview_img):
    slide.shapes.add_picture(overview_img, Inches(6.8), Inches(1.0), Inches(5.2), Inches(5.2))

add_shape(slide, Inches(0), Inches(6.5), W, Inches(0.005), RGBColor(0x33, 0x33, 0x55))
t3 = tb(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.6))
set_text(t3.text_frame, "Product Overview  |  Concept Stage  |  2026", size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 2: Architecture — Three-Tier System
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t.text_frame, "Three-Tier Architecture", size=36, color=WHITE, bold=True)

# Tier boxes
tiers = [
    ("SENSOR NODE", "1 per room", ACCENT_TEAL,
     ["nRF52840 SoC (BLE 5.3)", "CO2 (SCD41) + T/H (SHT40) + PM (PMSA003I)",
      "2x AA batteries, >12 month target", "Deep sleep 99% of the time",
      "BLE advertisement every 5 min"]),
    ("GATEWAY", "1 per floor", ACCENT_BLUE,
     ["ESP32-S3 + Ethernet", "Passively scans BLE advertisements",
      "Aggregates data from up to 50 nodes", "Forwards to cloud via MQTT/TLS",
      "Mains-powered, always on"]),
    ("CLOUD BACKEND", "Centralized", ACCENT_PURPLE,
     ["MQTT ingestion (AWS IoT Core)", "TimescaleDB time-series storage",
      "REST API + web dashboard", "Alert engine (CO2, temp, device health)",
      "OTA firmware deployment to fleet"]),
]

for i, (name, subtitle, color, bullets) in enumerate(tiers):
    x = Inches(0.6) + Inches(4.1) * i
    add_shape(slide, x, Inches(1.6), Inches(3.8), Inches(0.06), color)
    add_shape(slide, x, Inches(1.66), Inches(3.8), Inches(4.5), CARD_BG)

    t = tb(slide, x + Inches(0.25), Inches(1.85), Inches(3.3), Inches(0.5))
    set_text(t.text_frame, name, size=20, color=color, bold=True)

    t = tb(slide, x + Inches(0.25), Inches(2.35), Inches(3.3), Inches(0.3))
    set_text(t.text_frame, subtitle, size=12, color=LIGHT_GRAY, bold=True)

    for j, bullet in enumerate(bullets):
        t = tb(slide, x + Inches(0.25), Inches(2.85) + Inches(0.5) * j, Inches(3.3), Inches(0.4))
        set_text(t.text_frame, bullet, size=12, color=SOFT_WHITE)

# Arrows between tiers
for i in range(2):
    ax = Inches(4.4) + Inches(4.1) * i
    ay = Inches(3.8)
    add_shape(slide, ax, ay, Inches(0.3), Inches(0.04), LIGHT_GRAY)
    # Arrow label
    labels = ["BLE 5.3", "MQTT/TLS"]
    t = tb(slide, ax - Inches(0.1), ay - Inches(0.35), Inches(0.6), Inches(0.3))
    set_text(t.text_frame, labels[i], size=9, color=LIGHT_GRAY, bold=True, align=PP_ALIGN.CENTER)

# Bottom — data flow
add_shape(slide, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.7), CARD_BG)
t = tb(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5))
set_text(t.text_frame,
         "Sensor node (BLE adv, one-way) --> Gateway (passive scan, Ethernet) --> Cloud (MQTT) --> Dashboard (REST API)",
         size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 3: Hardware Cross-Sections
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_TEAL)

t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t.text_frame, "Inside the Hardware", size=36, color=WHITE, bold=True)

# Left — sensor node cross-section
node_img = os.path.join(_DIR, "Cross-Section — Sensor Node.png")
if os.path.exists(node_img):
    slide.shapes.add_picture(node_img, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))

add_shape(slide, Inches(0.5), Inches(1.15), Inches(5.5), Inches(0.35), ACCENT_TEAL)
t = tb(slide, Inches(0.65), Inches(1.18), Inches(5.2), Inches(0.3))
set_text(t.text_frame, "SENSOR NODE  (1 per room)", size=13, color=WHITE, bold=True)

# Right — gateway cross-section
gw_img = os.path.join(_DIR, "Cross-Section — Gateway.png")
if os.path.exists(gw_img):
    slide.shapes.add_picture(gw_img, Inches(6.5), Inches(1.5), Inches(6.3), Inches(3.9))

add_shape(slide, Inches(6.5), Inches(1.15), Inches(6.3), Inches(0.35), ACCENT_BLUE)
t = tb(slide, Inches(6.65), Inches(1.18), Inches(6.0), Inches(0.3))
set_text(t.text_frame, "GATEWAY  (1 per floor)", size=13, color=WHITE, bold=True)

# Key specs below gateway image
gw_specs = [
    ("MCU", "ESP32-S3 (BLE + WiFi + Ethernet MAC)"),
    ("Uplink", "100 Mbps Ethernet, MQTT over TLS"),
    ("Power", "Mains-powered via USB-C, always on"),
    ("Capacity", "Scans up to 50 BLE sensor nodes"),
]

for i, (label, value) in enumerate(gw_specs):
    y = Inches(5.6) + Inches(0.38) * i
    bg = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    add_shape(slide, Inches(6.5), y, Inches(6.3), Inches(0.32), bg)

    t = tb(slide, Inches(6.65), y + Inches(0.04), Inches(1.2), Inches(0.22))
    set_text(t.text_frame, label, size=10, color=ACCENT_BLUE, bold=True)

    t = tb(slide, Inches(7.9), y + Inches(0.04), Inches(4.7), Inches(0.22))
    set_text(t.text_frame, value, size=10, color=SOFT_WHITE)

# ============================================================
# SLIDE 4: Sensor Node Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_TEAL)

t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t.text_frame, "Sensor Node — Components & Power", size=36, color=WHITE, bold=True)

# Left — sensors
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.5), ACCENT_TEAL)
t = tb(slide, Inches(1.0), Inches(1.55), Inches(5.4), Inches(0.4))
set_text(t.text_frame, "SENSORS", size=14, color=WHITE, bold=True)

sensors = [
    ("SCD41 (Sensirion)", "CO2 ppm + temp + humidity", "+/-40 ppm accuracy", "$8-12", ACCENT_TEAL),
    ("SHT40 (Sensirion)", "Temperature + humidity", "+/-0.2C, +/-1.8% RH", "$1-2", ACCENT_TEAL),
    ("PMSA003I (Plantower)", "PM1.0, PM2.5, PM10", "+/-10 ug/m3", "$7-9", ACCENT_ORANGE),
]

for i, (name, measures, accuracy, cost, color) in enumerate(sensors):
    y = Inches(2.2) + Inches(0.95) * i
    bg = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    add_shape(slide, Inches(0.8), y, Inches(5.8), Inches(0.82), bg)
    accent_bar(slide, Inches(0.8), y + Inches(0.1), Inches(0.6), color)

    t = tb(slide, Inches(1.1), y + Inches(0.05), Inches(2.5), Inches(0.3))
    set_text(t.text_frame, name, size=13, color=WHITE, bold=True)

    t = tb(slide, Inches(3.8), y + Inches(0.05), Inches(2.5), Inches(0.3))
    set_text(t.text_frame, cost, size=13, color=color, bold=True, align=PP_ALIGN.RIGHT)

    t = tb(slide, Inches(1.1), y + Inches(0.35), Inches(3.0), Inches(0.3))
    set_text(t.text_frame, f"{measures}  |  {accuracy}", size=11, color=LIGHT_GRAY)

# Right — power & firmware
rx = Inches(7.3)
add_shape(slide, rx, Inches(1.5), Inches(5.3), Inches(0.5), ACCENT_ORANGE)
t = tb(slide, rx + Inches(0.2), Inches(1.55), Inches(4.8), Inches(0.4))
set_text(t.text_frame, "POWER ARCHITECTURE", size=14, color=WHITE, bold=True)

power_items = [
    ("Source", "2x AA lithium (Energizer L91)"),
    ("Capacity", "3000 mAh, 6V input -> TPS62740 -> 3.3V"),
    ("Sleep current", "~2 uA (MCU + regulator quiescent)"),
    ("Sense cycle", "5s active every 5 min = 99.7% sleeping"),
    ("Target life", ">12 months (PM sensor is the bottleneck)"),
    ("Battery issue", "PM sensor @ 25 mA eats 417 uA avg alone"),
]

for i, (label, value) in enumerate(power_items):
    y = Inches(2.2) + Inches(0.63) * i
    bg = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    add_shape(slide, rx, y, Inches(5.3), Inches(0.52), bg)

    t = tb(slide, rx + Inches(0.15), y + Inches(0.1), Inches(1.4), Inches(0.3))
    set_text(t.text_frame, label, size=11, color=ACCENT_ORANGE, bold=True)

    t = tb(slide, rx + Inches(1.6), y + Inches(0.1), Inches(3.5), Inches(0.3))
    set_text(t.text_frame, value, size=11, color=SOFT_WHITE)

# Bottom — firmware
add_shape(slide, Inches(0.8), Inches(5.4), Inches(11.8), Inches(0.5), ACCENT_BLUE)
t = tb(slide, Inches(1.0), Inches(5.45), Inches(11.4), Inches(0.4))
set_text(t.text_frame, "FIRMWARE (Zephyr RTOS on nRF52840)", size=14, color=WHITE, bold=True)

fw_modules = [
    ("Sensor Mgr", "Read SCD41, SHT40, PMSA003I\non schedule. Power-gate PM sensor."),
    ("BLE Advertiser", "Encode readings into 20-byte\nextended advertisement payload."),
    ("Power Mgr", "Deep sleep entry/exit,\nbattery ADC monitoring."),
    ("OTA Mgr", "BLE DFU via MCUboot.\nA/B partitioning, rollback on failure."),
    ("Config Mgr", "Sample interval, PM enable/disable.\nBLE write from gateway."),
]

for i, (name, desc) in enumerate(fw_modules):
    x = Inches(0.8) + Inches(2.46) * i
    add_shape(slide, x, Inches(6.05), Inches(2.26), Inches(1.1), CARD_BG)
    t = tb(slide, x + Inches(0.1), Inches(6.1), Inches(2.05), Inches(0.3))
    set_text(t.text_frame, name, size=11, color=ACCENT_BLUE, bold=True)
    t = tb(slide, x + Inches(0.1), Inches(6.4), Inches(2.05), Inches(0.6))
    set_text(t.text_frame, desc, size=9, color=LIGHT_GRAY)

# ============================================================
# SLIDE 4: Constraints & BOM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_ORANGE)

t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t.text_frame, "Constraints & Cost", size=36, color=WHITE, bold=True)

# Left — constraints
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.5), ACCENT_ORANGE)
t = tb(slide, Inches(1.0), Inches(1.55), Inches(5.4), Inches(0.4))
set_text(t.text_frame, "HARD CONSTRAINTS", size=14, color=WHITE, bold=True)

constraints = [
    ("Battery life", ">12 months on 2x AA", "Mount and forget. No wiring, no charging."),
    ("Node BOM", "<$35 at 1k units", "50-200 nodes per building -- cost must scale."),
    ("Gateway BOM", "<$40 at 1k units", "1 per floor. Justifies itself vs. WiFi."),
    ("CO2 accuracy", "+/-50 ppm + 5%", "Facility managers make HVAC decisions from this."),
    ("Certification", "FCC, CE, IC", "BLE = intentional radiator. All target markets."),
    ("Density", "50 nodes per gateway", "BLE scanning must handle without packet loss."),
]

for i, (name, value, note) in enumerate(constraints):
    y = Inches(2.2) + Inches(0.75) * i
    bg = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    add_shape(slide, Inches(0.8), y, Inches(5.8), Inches(0.63), bg)

    t = tb(slide, Inches(1.0), y + Inches(0.05), Inches(2.0), Inches(0.3))
    set_text(t.text_frame, name, size=12, color=ACCENT_ORANGE, bold=True)

    t = tb(slide, Inches(1.0), y + Inches(0.05), Inches(5.4), Inches(0.3))
    set_text(t.text_frame, value, size=12, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

    t = tb(slide, Inches(1.0), y + Inches(0.33), Inches(5.4), Inches(0.25))
    set_text(t.text_frame, note, size=10, color=LIGHT_GRAY)

# Right — BOM breakdown
rx = Inches(7.3)
add_shape(slide, rx, Inches(1.5), Inches(5.3), Inches(0.5), ACCENT_TEAL)
t = tb(slide, rx + Inches(0.2), Inches(1.55), Inches(4.8), Inches(0.4))
set_text(t.text_frame, "SENSOR NODE BOM (1k units)", size=14, color=WHITE, bold=True)

bom = [
    ("nRF52840 SoC", "$4.50"),
    ("SCD41 CO2 sensor", "$10.00"),
    ("SHT40 temp/humidity", "$1.50"),
    ("PMSA003I PM sensor", "$8.00"),
    ("TPS62740 regulator", "$1.50"),
    ("RGB LED + passives", "$0.50"),
    ("PCB (45x45mm, 4-layer)", "$1.50"),
    ("2x AA battery holder", "$0.40"),
    ("2x AA lithium cells", "$2.50"),
    ("Enclosure (ABS molded)", "$2.50"),
    ("Assembly + test", "$2.50"),
]

for i, (item, cost) in enumerate(bom):
    y = Inches(2.15) + Inches(0.35) * i
    bg = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    add_shape(slide, rx, y, Inches(5.3), Inches(0.28), bg)

    t = tb(slide, rx + Inches(0.15), y + Inches(0.02), Inches(3.2), Inches(0.22))
    set_text(t.text_frame, item, size=10, color=SOFT_WHITE)

    t = tb(slide, rx + Inches(3.5), y + Inches(0.02), Inches(1.5), Inches(0.22))
    set_text(t.text_frame, cost, size=10, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

# Total
total_y = Inches(2.15) + Inches(0.35) * len(bom) + Inches(0.08)
add_shape(slide, rx, total_y, Inches(5.3), Inches(0.35), ACCENT_TEAL)
t = tb(slide, rx + Inches(0.15), total_y + Inches(0.04), Inches(3.2), Inches(0.25))
set_text(t.text_frame, "Total COGS (sensor node)", size=11, color=WHITE, bold=True)
t = tb(slide, rx + Inches(3.5), total_y + Inches(0.04), Inches(1.5), Inches(0.25))
set_text(t.text_frame, "~$35.90", size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

# Gateway BOM summary
gy = total_y + Inches(0.55)
add_shape(slide, rx, gy, Inches(5.3), Inches(0.5), CARD_BG)
t = tb(slide, rx + Inches(0.15), gy + Inches(0.05), Inches(3.5), Inches(0.2))
set_text(t.text_frame, "Gateway: ESP32-S3 + Ethernet PHY + PSU + enclosure", size=10, color=LIGHT_GRAY)
t = tb(slide, rx + Inches(3.5), gy + Inches(0.05), Inches(1.5), Inches(0.2))
set_text(t.text_frame, "~$35", size=10, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
t = tb(slide, rx + Inches(0.15), gy + Inches(0.28), Inches(4.8), Inches(0.2))
set_text(t.text_frame, "Cloud hosting: <$0.50/node/month", size=10, color=LIGHT_GRAY)

# ============================================================
# SLIDE 5: Hardest Problems & Component Tradeoffs
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_RED)

t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t.text_frame, "Hardest Problems & Component Tradeoffs", size=32, color=WHITE, bold=True)

# Left — three hardest problems
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.5), ACCENT_RED)
t = tb(slide, Inches(1.0), Inches(1.55), Inches(5.4), Inches(0.4))
set_text(t.text_frame, "THREE HARDEST PROBLEMS", size=14, color=WHITE, bold=True)

problems = [
    ("PM sensor vs. battery life",
     "The PMSA003I draws 25 mA during sampling -- 50x more than everything else combined. "
     "At 5 min intervals, it alone exceeds the entire power budget. "
     "Options: reduce PM sampling to 30-60 min, make PM a removable module, or accept shorter battery."),
    ("BLE range in concrete offices",
     "BLE advertising reach drops to 5-10m through concrete walls. "
     "1 gateway per floor assumes 15m line-of-sight. "
     "Dense offices may need 2-3 gateways per floor, increasing cost."),
    ("OTA over BLE to sleeping nodes",
     "Firmware updates via cloud -> gateway -> BLE DFU to a node "
     "that sleeps 99% of the time. Must coordinate wake windows, "
     "handle interrupted transfers, and never brick deployed devices."),
]

for i, (title, desc) in enumerate(problems):
    y = Inches(2.2) + Inches(1.35) * i
    add_shape(slide, Inches(0.8), y, Inches(5.8), Inches(1.15), CARD_BG)
    circle_num(slide, Inches(1.0), y + Inches(0.15), i + 1, ACCENT_RED)

    t = tb(slide, Inches(1.6), y + Inches(0.08), Inches(4.8), Inches(0.35))
    set_text(t.text_frame, title, size=15, color=WHITE, bold=True)

    t = tb(slide, Inches(1.6), y + Inches(0.42), Inches(4.8), Inches(0.7))
    set_text(t.text_frame, desc, size=11, color=LIGHT_GRAY)

# Right — component tradeoffs
rx = Inches(7.3)
add_shape(slide, rx, Inches(1.5), Inches(5.3), Inches(0.5), ACCENT_TEAL)
t = tb(slide, rx + Inches(0.2), Inches(1.55), Inches(4.8), Inches(0.4))
set_text(t.text_frame, "COMPONENT TRADEOFFS", size=14, color=WHITE, bold=True)

tradeoffs = [
    ("nRF52840", "Performance",
     "Best BLE sleep ($4.50) vs. ESP32-C3 ($1.50, worse sleep). 12-month battery forces the premium."),
    ("SCD41 CO2", "Performance",
     "Only photoacoustic CO2 at this size + accuracy. $10 dominates BOM. No alternative meets spec."),
    ("PMSA003I PM", "Cost vs. power",
     "Adds $8 + destroys battery budget. Removable module (two SKUs) lets the customer decide."),
    ("TPS62740", "Power efficiency",
     "360 nA quiescent ($1.50) vs. LDO ($0.30, wastes months of battery in quiescent alone)."),
    ("ESP32-S3 (GW)", "Cost",
     "Cheapest BLE+Ethernet path ($3-4). Mains-powered -- no power tension at all."),
]

for i, (component, axis, desc) in enumerate(tradeoffs):
    y = Inches(2.2) + Inches(0.88) * i
    bg = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    add_shape(slide, rx, y, Inches(5.3), Inches(0.76), bg)

    t = tb(slide, rx + Inches(0.2), y + Inches(0.05), Inches(2.2), Inches(0.3))
    set_text(t.text_frame, component, size=13, color=WHITE, bold=True)

    t = tb(slide, rx + Inches(2.6), y + Inches(0.05), Inches(2.5), Inches(0.3))
    set_text(t.text_frame, axis, size=12, color=ACCENT_TEAL, bold=True, align=PP_ALIGN.RIGHT)

    t = tb(slide, rx + Inches(0.2), y + Inches(0.33), Inches(4.9), Inches(0.4))
    set_text(t.text_frame, desc, size=10, color=LIGHT_GRAY)

# Bottom
t = tb(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6))
set_text(t.text_frame,
         "The PM sensor is the central tension: it differentiates the product but threatens the battery target. Two-SKU model is the likely resolution.",
         size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 6: Key Decisions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_PURPLE)

t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t.text_frame, "Key Technical Decisions", size=36, color=WHITE, bold=True)

decisions = [
    ("BLE + Gateway vs. WiFi Direct", "CHOSEN: BLE + GATEWAY",
     "WiFi draws 100-300 mA during Tx vs. 8 mA for BLE advertising. WiFi would drain AAs in weeks. "
     "The gateway adds $35-40 per floor but enables >12 month battery life per sensor node.",
     "Adds a second hardware product (gateway). BLE range limits placement to ~15m LOS. "
     "1 gateway per ~500 m2.",
     ACCENT_TEAL),
    ("AA Batteries vs. Rechargeable LiPo", "CHOSEN: 2x AA LITHIUM",
     "Facility staff can swap AAs in seconds with no tools. Rechargeable adds USB-C, charge IC, and "
     "'another thing to charge' -- doesn't fit mount-and-forget deployment.",
     "Non-rechargeable: ~$6/year per node in batteries. 200-node building = $1,200/year. "
     "Rechargeable V2 SKU possible if this becomes a concern.",
     ACCENT_ORANGE),
    ("Cloud-Only vs. Edge Analytics", "CHOSEN: CLOUD-ONLY (V1)",
     "Sensor nodes sleep 99%. Gateway could do edge analytics, but adds FW complexity for marginal "
     "latency gain -- dashboard refreshes every 30s anyway. Cloud processing keeps FW simple and allows "
     "algorithm updates without OTA.",
     "Requires reliable Ethernet. Alert latency is cloud-round-trip (~1-5s). "
     "Acceptable for HVAC decisions (minutes-scale).",
     ACCENT_PURPLE),
]

for i, (title, choice, rationale, consequences, color) in enumerate(decisions):
    y = Inches(1.5) + Inches(1.8) * i
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.6), CARD_BG)
    add_shape(slide, Inches(0.8), y, Inches(0.1), Inches(1.6), color)

    circle_num(slide, Inches(1.1), y + Inches(0.15), i + 1, color)

    t = tb(slide, Inches(1.7), y + Inches(0.1), Inches(4.0), Inches(0.3))
    set_text(t.text_frame, title, size=16, color=WHITE, bold=True)

    t = tb(slide, Inches(6.0), y + Inches(0.1), Inches(3.0), Inches(0.3))
    set_text(t.text_frame, choice, size=12, color=color, bold=True)

    t = tb(slide, Inches(1.7), y + Inches(0.5), Inches(5.0), Inches(0.8))
    set_text(t.text_frame, rationale, size=11, color=SOFT_WHITE)

    t = tb(slide, Inches(7.0), y + Inches(0.5), Inches(5.2), Inches(0.8))
    tf = t.text_frame
    set_text(tf, "Consequences:", size=10, color=ACCENT_ORANGE, bold=True)
    add_p(tf, consequences, size=10, color=LIGHT_GRAY, before=Pt(2))

# ============================================================
# SLIDE 7: Open Questions & Risks
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_ORANGE)

t = tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(t.text_frame, "Open Questions & Risks", size=36, color=WHITE, bold=True)

# Header row
add_shape(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.5), ACCENT_ORANGE)
headers = [
    (Inches(0.9), Inches(0.4), "#"),
    (Inches(1.4), Inches(5.0), "Question / Risk"),
    (Inches(6.6), Inches(0.8), "Impact"),
    (Inches(7.5), Inches(1.0), "Target"),
]
for x, w, label in headers:
    t = tb(slide, x, Inches(1.45), w, Inches(0.35))
    set_text(t.text_frame, label, size=11, color=WHITE, bold=True)

questions = [
    ("1", "PM sensor power budget exceeds target -- reduce sample rate, make removable module, or accept shorter battery?", "H", "M2"),
    ("2", "BLE range in concrete offices -- need field testing in 3+ real buildings", "M", "M3"),
    ("3", "SCD41 warm-up: does 5s window maintain +/-50 ppm? If 15s needed, CO2 power 3x.", "H", "M2"),
    ("4", "OTA reliability over BLE DFU to sleeping nodes -- retry strategy needed", "M", "M3"),
    ("5", "SCD41 sole source risk -- Sensirion only, 12-26 week lead times historically", "H", "M1"),
    ("6", "Gateway density: 1/floor vs. 1/zone. Need field testing.", "M", "M4"),
    ("7", "Battery replacement economics: $1,200/yr for 200-node building. Acceptable?", "L", "M6"),
]

for i, (num, question, impact, target) in enumerate(questions):
    y = Inches(2.05) + Inches(0.65) * i
    bg = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.55), bg)

    impact_color = ACCENT_RED if impact == "H" else (ACCENT_ORANGE if impact == "M" else LIGHT_GRAY)

    t = tb(slide, Inches(0.9), y + Inches(0.1), Inches(0.4), Inches(0.3))
    set_text(t.text_frame, num, size=11, color=LIGHT_GRAY, bold=True)

    t = tb(slide, Inches(1.4), y + Inches(0.1), Inches(5.0), Inches(0.35))
    set_text(t.text_frame, question, size=11, color=SOFT_WHITE)

    t = tb(slide, Inches(6.6), y + Inches(0.1), Inches(0.6), Inches(0.3))
    set_text(t.text_frame, impact, size=12, color=impact_color, bold=True, align=PP_ALIGN.CENTER)

    t = tb(slide, Inches(7.5), y + Inches(0.1), Inches(0.8), Inches(0.3))
    set_text(t.text_frame, target, size=11, color=LIGHT_GRAY, bold=True)

# Bottom note
add_shape(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.45), CARD_BG)
t = tb(slide, Inches(1.0), Inches(6.75), Inches(11.3), Inches(0.3))
set_text(t.text_frame,
         "Items #1 and #3 are the highest-priority: both threaten the 12-month battery target, which is the product's core promise.",
         size=12, color=ACCENT_ORANGE)

# ============================================================
# Save
# ============================================================
output_path = os.path.join(_DIR, "AirSense_Deck.pptx")
prs.save(output_path)
print(f"Saved to {output_path}")
